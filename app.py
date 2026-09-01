import json
import re
import tempfile
from io import BytesIO
from datetime import date
from pathlib import Path
from typing import Any, Dict, List

import streamlit as st
from docx import Document
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from pptx.dml.color import RGBColor

APP_TITLE = "Startdocument Generator"
APP_DIR = Path(__file__).resolve().parent
TEMPLATE_CANDIDATES = [
    APP_DIR / "templates" / "Startdocument_Cooble_template.pptx",
    APP_DIR / "Startdocument_Cooble_template.pptx",
    Path("templates/Startdocument_Cooble_template.pptx"),
    Path("Startdocument_Cooble_template.pptx"),
]

def get_template_path() -> Path:
    for candidate in TEMPLATE_CANDIDATES:
        if candidate.exists():
            return candidate
    checked = ", ".join(str(c) for c in TEMPLATE_CANDIDATES)
    raise RuntimeError(
        "PowerPoint-template niet gevonden. Controleer of 'Startdocument_Cooble_template.pptx' "
        "in de map 'templates' staat of in de hoofdmap van de repository. Gecheckte paden: " + checked
    )
DEFAULT_MODEL = "gpt-4.1"


# v2.1: stabiele demografie + presentatie-engine dichter op Cooble-template; grotere leesbare fonts.
def generate_with_openai_pipeline(vacature: str, intake: str, linkedin_size: str, extra: str, status=None) -> Dict[str, Any]:
    if status:
        status.write("Stap 1/7: feiten uit vacature en intake halen")
    facts = call_openai_json(build_fact_extraction_prompt(vacature, intake, extra), use_web=False)

    fallback = extract_basis_fallback(vacature, intake)
    for fk in ["klantnaam", "vacaturenaam", "salaris"]:
        if is_empty_or_placeholder(facts.get(fk, "")) and fallback.get(fk):
            facts[fk] = fallback[fk]
    facts["salaris"] = normalize_salary_display(facts.get("salaris", ""))

    extracted_no_go = extract_no_go_companies_from_intake(intake + "\n" + extra)
    if extracted_no_go:
        merged = []
        for item in clean_list(facts.get("no_go_bedrijven", [])) + extracted_no_go:
            c = clean_company_name(item)
            if c and c not in merged:
                merged.append(c)
        facts["no_go_bedrijven"] = merged

    if status:
        status.write("Stap 2/8: doelgroep en concurrenten online onderzoeken")
    market = call_openai_json(build_target_market_research_prompt(facts, linkedin_size), use_web=True)

    if status:
        status.write("Stap 3/8: belangrijkste arbeidsvoorwaarden online onderzoeken")
    conditions = call_openai_json(build_employment_conditions_research_prompt(facts), use_web=True)

    if status:
        status.write("Stap 4/8: pullfactoren online onderzoeken")
    pull = call_openai_json(build_pullfactors_research_prompt(facts), use_web=True)
    if pullfactors_are_invalid(pull.get("pullfactoren", []), facts.get("klantnaam", "")):
        pull = call_openai_json(build_pullfactors_research_prompt(facts, strict_retry=True), use_web=True)
    if pullfactors_are_invalid(pull.get("pullfactoren", []), facts.get("klantnaam", "")):
        raise RuntimeError("Online onderzoek leverde geen geldige pullfactoren op. De tool stopt bewust in plaats van vacature-inhoud of arbeidsmarktkrapte als pullfactor te gebruiken.")

    if status:
        status.write("Stap 5/8: leeftijd en man-vrouwverhouding online onderzoeken")
    demographics = call_openai_json(build_demographics_research_prompt(facts), use_web=True)
    research = merge_research_parts(market, conditions, pull, demographics)

    if status:
        status.write("Stap 6/8: startdocument-content schrijven")
    data = call_openai_json(build_writer_prompt(facts, research, vacature, intake, linkedin_size, extra), use_web=False)

    if status:
        status.write("Stap 7/8: presentatiekwaliteit aanscherpen")
    try:
        data = call_openai_json(build_presentation_prompt(data, facts, research), use_web=False)
    except Exception as presentation_error:
        data.setdefault("kwaliteitscontrole", {}).setdefault("waarschuwingen", []).append(str(presentation_error))

    if status:
        status.write("Stap 8/8: business rules toepassen")
    data = apply_business_rules(data, intake + "\n" + extra, linkedin_size, vacature, extra)
    # Researchvelden zijn leidend voor externe marktdata.
    data.setdefault("doelgroepanalyse", {})["pullfactoren"] = presentation_bullets(normalize_pullfactors(research.get("pullfactoren", [])), 3)
    data.setdefault("voorwaarden", {})["belangrijkste_arbeidsvoorwaarden"] = presentation_bullets(normalize_conditions(research.get("belangrijkste_arbeidsvoorwaarden", [])), 3)
    stable_demo = deterministic_demographics(facts, research)
    data.setdefault("doelgroepanalyse", {})["geslacht"] = stable_demo.get("geslacht", {"man": "", "vrouw": ""})
    data.setdefault("doelgroepanalyse", {})["leeftijdsverdeling"] = stable_demo.get("leeftijdsverdeling", normalize_age_distribution(research.get("leeftijdsverdeling", [])))
    data.setdefault("basisgegevens", {})["salaris"] = normalize_salary_display(data.get("basisgegevens", {}).get("salaris", ""))
    data.setdefault("kwaliteitscontrole", {})["pipeline"] = "v2.1: facts -> occupation-only web market -> web conditions -> validated web pull -> deterministic demographics -> writer -> template-first presentation"
    return data


st.set_page_config(page_title=APP_TITLE, page_icon="📄", layout="wide")


def read_docx(file) -> str:
    doc = Document(file)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    table_text = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_text.append(" | ".join(cells))
    return "\n".join(paragraphs + table_text)


def read_pdf(file) -> str:
    reader = PdfReader(file)
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)
    return "\n".join(pages)


def read_uploaded_file(file) -> str:
    if file is None:
        return ""
    name = file.name.lower()
    try:
        if name.endswith(".docx"):
            return read_docx(file)
        if name.endswith(".pdf"):
            return read_pdf(file)
        if name.endswith(".txt"):
            return file.read().decode("utf-8", errors="ignore")
    except Exception as exc:
        raise RuntimeError(f"Kon bestand '{file.name}' niet uitlezen: {exc}")
    return ""


def clean_list(items: List[str]) -> List[str]:
    return [str(x).strip(" •-\n\t") for x in (items or []) if str(x).strip(" •-\n\t")]


def split_one_topic_per_bullet(items: List[str]) -> List[str]:
    """Zorgt dat samengestelde bullets worden opgeknipt naar één onderwerp per bullet."""
    result: List[str] = []
    for raw in clean_list(items):
        text = raw.strip()
        # Splits alleen op duidelijke opsommingen, niet op woorden als "Learning & Development".
        parts = re.split(r"\s*(?:;|/|\+|,\s*(?=[A-ZÀ-ÖØ-Þ]))\s*", text)
        expanded: List[str] = []
        for part in parts:
            # Splits op " en " alleen als beide delen kort genoeg zijn om losse thema's te zijn.
            sub = re.split(r"\s+en\s+", part)
            if len(sub) == 2 and all(1 <= len(x.split()) <= 4 for x in sub):
                expanded.extend(sub)
            else:
                expanded.append(part)
        for part in expanded:
            part = part.strip(" •-\n\t")
            if part and part not in result:
                result.append(part)
    return result




def first_nonempty_line(text: str) -> str:
    for line in (text or "").splitlines():
        line = line.strip()
        if line and len(line) < 90:
            return line
    return ""


def extract_basis_fallback(vacature_text: str, intake_text: str) -> Dict[str, str]:
    """Deterministische fallback voor klantnaam, vacaturenaam en salaris uit input."""
    combined = f"{intake_text}\n{vacature_text}"
    fallback: Dict[str, str] = {}

    # Klantnaam uit intakevelden of herkenbare vacaturetekst.
    m = re.search(r"(?im)^\s*Klant\s*:\s*(.+)$", combined)
    if m:
        fallback["klantnaam"] = m.group(1).strip()
    else:
        m = re.search(r"(?i)\bWij zijn\s+([A-Z][A-Za-z0-9&+ .'-]{2,40})", vacature_text)
        if m:
            name = re.split(r"[\n\.;,]", m.group(1).strip())[0].strip()
            fallback["klantnaam"] = name

    # Vacaturenaam uit intakeveld; als leeg, eerste regel vacaturetekst.
    m = re.search(r"(?im)^\s*Vacature\s*:\s*(.+)$", combined)
    if m and m.group(1).strip():
        fallback["vacaturenaam"] = m.group(1).strip()
    else:
        title = first_nonempty_line(vacature_text)
        # Vermijd plaats/metadata als titel.
        if title and not re.search(r"\b(remote|hybrid|nl|ov|sp|locatie)\b", title, re.I):
            fallback["vacaturenaam"] = title

    # Salaris: euro-range, salarisschaal/schaal of expliciete salarisregel.
    salary_patterns = [
        r"(?i)(?:salaris(?:range)?|salarisindicatie)\s*[:\-]?\s*([^\n]{3,80})",
        r"(?i)\b(schaal\s*[0-9][0-9A-Za-z/\- ]{0,30})",
        r"(?i)(€\s?[\d\.,]+\s*(?:-|–|tot)\s*€?\s?[\d\.,]+)",
        r"(?i)([\d\.,]+\s*(?:-|–|tot)\s*[\d\.,]+\s*(?:euro|bruto|per maand)?)",
    ]
    for pat in salary_patterns:
        m = re.search(pat, combined)
        if m:
            val = m.group(1).strip(" .;,")
            if val and not re.search(r"weten we niet|onbekend", val, re.I):
                fallback["salaris"] = val
                break
            # Als er staat: salaris weten we niet, schaal 8/9/10, pak alsnog schaal.
            scale = re.search(r"(?i)\b(schaal\s*[0-9][0-9A-Za-z/\- ]{0,30})", m.group(0))
            if scale:
                fallback["salaris"] = scale.group(1).strip()
                break
    return fallback


def is_empty_or_placeholder(value: str) -> bool:
    value = str(value or "").strip()
    return value == "" or value.lower() in {"onbekend", "n.v.t.", "nvt", "in overleg", "-", "..."}


def is_placeholder_company(value: str) -> bool:
    return bool(re.search(r"(?i)^bedrijf\s+[a-z0-9]$|^concurrent\s+[a-z0-9]$|^organisatie\s+[a-z0-9]$", str(value or "").strip()))


def infer_competitors_offline(vacature_text: str, intake_text: str, current: List[str]) -> List[str]:
    """Fallback wanneer de AI placeholders geeft. Beperkt, maar beter dan Bedrijf A/B/C."""
    text = f"{vacature_text}\n{intake_text}".lower()
    companies = []
    # Bedrijven die expliciet in intake staan, mogen ook als concurrent/check-eerst zichtbaar worden.
    companies.extend(extract_no_go_companies_from_intake(intake_text))
    maps = [
        (["waterkwaliteit", "afvalwater", "waterwet", "waterschap", "watermanagement"], ["Witteveen+Bos", "Royal HaskoningDHV", "Sweco", "Antea Group", "Arcadis", "TAUW"]),
        (["luchtkwaliteit", "milieuconsultant", "emissie", "vergunning", "omgevingswet"], ["Royal HaskoningDHV", "Witteveen+Bos", "Sweco", "Antea Group", "Arcadis", "TAUW"]),
        (["business analist", "informatieanalist", "product owner"], ["Sogeti", "Capgemini", "Ordina", "CGI", "Conclusion", "Atos"]),
        (["lead engineer", "engineer", "warmtenet", "ondergrondse infra"], ["BAM", "Heijmans", "VolkerWessels", "Strukton", "Equans", "SPIE"]),
        (["accountmanager", "retentie", "upsell", "sales manager"], ["LeasePlan", "Alphabet", "Arval", "Athlon", "ALD Automotive"]),
    ]
    for keywords, names in maps:
        if any(k in text for k in keywords):
            companies.extend(names)
            break
    for item in current or []:
        if item and not is_placeholder_company(item):
            companies.append(item)
    result = []
    for c in companies:
        c = clean_company_name(c)
        if c and not is_placeholder_company(c) and c not in result:
            result.append(c)
    return result[:8]


ALLOWED_EMPLOYMENT_CONDITIONS = [
    "Salaris",
    "Pensioenregeling",
    "Vakantiedagen",
    "ADV-dagen",
    "Bonusregeling",
    "Eindejaarsuitkering",
    "13e maand",
    "Leaseauto",
    "Mobiliteitsbudget",
    "Reiskostenvergoeding",
    "Thuiswerkregeling",
    "Thuiswerkvergoeding",
    "Flexibele werktijden",
    "Opleidingsbudget",
    "Studiekostenregeling",
    "Ploegentoeslag",
    "Overwerkvergoeding",
    "Onkostenvergoeding",
    "Aandelenregeling",
]


def normalize_condition_label(text: str) -> str:
    """Map webresearch uitsluitend naar echte primaire/secundaire arbeidsvoorwaarden.

    Cultuur, inhoud, autonomie, ontwikkelkansen en andere pullfactoren worden bewust
    niet toegelaten in dit veld.
    """
    text = str(text or "").strip(" •-\n\t.,;:")
    low = re.sub(r"\s+", " ", text.lower()).strip()
    if not low:
        return ""

    # Expliciet blokkeren: dit zijn geen primaire/secundaire arbeidsvoorwaarden.
    blocked = (
        "werksfeer", "cultuur", "autonomie", "impact", "inhoud",
        "uitdaging", "doorgroei", "groeimogelijk", "ontwikkelmogelijk",
        "professionele ontwikkeling", "werkzekerheid", "baanzekerheid",
        "maatschappelijke relevantie", "verantwoordelijkheid", "team",
        "leidinggevende", "zingeving", "erkenning", "carriere", "carrière"
    )
    if any(term in low for term in blocked):
        return ""

    rules = [
        (("salaris", "loon", "basissalaris", "beloning", "salary", "pay"), "Salaris"),
        (("pensioen",), "Pensioenregeling"),
        (("adv", "atv"), "ADV-dagen"),
        (("vakantie", "vakantiedag", "verlofdagen", "vrije dagen"), "Vakantiedagen"),
        (("bonus", "variabele beloning", "prestatiebeloning"), "Bonusregeling"),
        (("eindejaarsuitkering", "eindejaarsbonus"), "Eindejaarsuitkering"),
        (("13e maand", "dertiende maand"), "13e maand"),
        (("leaseauto", "auto van de zaak", "bedrijfsauto"), "Leaseauto"),
        (("mobiliteitsbudget",), "Mobiliteitsbudget"),
        (("reiskosten", "kilometervergoeding", "ov-vergoeding"), "Reiskostenvergoeding"),
        (("thuiswerkregeling", "hybride werken", "thuiswerken", "remote werken"), "Thuiswerkregeling"),
        (("thuiswerkvergoeding",), "Thuiswerkvergoeding"),
        (("flexibele werktijden", "flextijden", "flexibele uren"), "Flexibele werktijden"),
        (("opleidingsbudget", "persoonlijk opleidingsbudget", "trainingsbudget"), "Opleidingsbudget"),
        (("studiekosten", "studiekostenregeling", "studievergoeding"), "Studiekostenregeling"),
        (("ploegentoeslag", "onregelmatigheidstoeslag"), "Ploegentoeslag"),
        (("overwerkvergoeding", "overwerktoeslag"), "Overwerkvergoeding"),
        (("onkostenvergoeding",), "Onkostenvergoeding"),
        (("aandelenregeling", "aandelenplan", "stock options"), "Aandelenregeling"),
    ]
    for needles, label in rules:
        if any(n in low for n in needles):
            return label

    # Alleen labels die letterlijk op de whitelist staan mogen door.
    for allowed in ALLOWED_EMPLOYMENT_CONDITIONS:
        if low == allowed.lower():
            return allowed
    return ""


def normalize_conditions(items: List[str]) -> List[str]:
    out: List[str] = []
    for item in clean_list(items or []):
        label = normalize_condition_label(item)
        if label and label not in out:
            out.append(label)
    return out[:3]


def employment_conditions_are_invalid(items: List[str]) -> bool:
    cleaned = normalize_conditions(items)
    return len(cleaned) != 3 or any(item not in ALLOWED_EMPLOYMENT_CONDITIONS for item in cleaned)

def normalize_pullfactor_label(text: str) -> str:
    """Maak pullfactoren compact, natuurlijk en presentatiewaardig zonder losse/rare fragmenten."""
    text = str(text or "").strip(" •-\n\t.,;:")
    text = re.sub(r"\s+", " ", text)
    if not text:
        return ""
    # Vang te losse éénwoord-items af die zonder context vreemd ogen.
    replacements = {
        "certificering": "Erkende certificeringen",
        "ontwikkeling": "Professionele ontwikkeling",
        "autonomie": "Meer autonomie",
        "doorgroei": "Doorgroeimogelijkheden",
        "flexibiliteit": "Meer flexibiliteit",
        "impact": "Zichtbare impact",
        "uitdaging": "Inhoudelijke uitdaging",
        "zekerheid": "Baanzekerheid",
    }
    low = text.lower()
    if low in replacements:
        text = replacements[low]
    # Pullfactoren moeten korte labels of korte natuurlijke zinnen zijn.
    words = text.split()
    if len(words) > 9:
        text = " ".join(words[:9]).rstrip(" ,;:-")
    return text[:1].upper() + text[1:] if text else text


def normalize_pullfactors(items: List[str]) -> List[str]:
    out: List[str] = []
    for item in clean_list(items):
        item = normalize_pullfactor_label(item)
        if item and item not in out:
            out.append(item)
    return out[:3]

def normalize_salary_display(value: str) -> str:
    """Toon salaris als alleen getal/range; bij salarisschalen mag het woord 'Schaal' blijven staan."""
    text = str(value or "").strip()
    if not text:
        return ""
    # Schaalnotatie: "Schaal 8/9/10" -> "Schaal 8/9/10"
    m = re.search(r"(?i)\bschaal\s*([0-9]+(?:\s*/\s*[0-9]+)*)", text)
    if m:
        return "Schaal " + re.sub(r"\s+", "", m.group(1))
    # Range met bedragen: € 5.200 - € 7.000 bruto p/m -> 5.200 - 7.000
    nums = re.findall(r"(?<![A-Za-z])(?:\d{1,3}(?:[. ]\d{3})+|\d{4,6})(?:,\d{1,2})?", text)
    if len(nums) >= 2:
        return f"{nums[0].replace(' ', '.')} - {nums[1].replace(' ', '.')}"
    if len(nums) == 1:
        return nums[0].replace(' ', '.')
    # Laat alleen cijfer/range-achtige tekens over als laatste veilige fallback.
    cleaned = re.sub(r"[^0-9.,/\-– ]", "", text)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" -–.,")
    return cleaned


def normalize_age_distribution(items: List[str]) -> List[str]:
    """Normaliseer uitsluitend extern onderzochte leeftijdsdata; verzin geen standaardverdeling."""
    values = clean_list(items)
    wanted = ["15-24", "25-34", "35-49", "50+"]
    parsed: Dict[str, int] = {}
    for val in values:
        text = str(val)
        pct_match = re.search(r"(\\d{1,3})\\s*%", text)
        if not pct_match:
            continue
        pct = max(0, min(100, int(pct_match.group(1))))
        for label in wanted:
            if label in text.replace(" ", "") or label in text:
                parsed[label] = pct
                break
    if len(parsed) != 4:
        return values[:4]
    total = sum(parsed.values())
    if total != 100:
        # Alleen afrondingsverschillen corrigeren; geen inhoudelijke herschatting.
        diff = 100 - total
        largest = max(parsed, key=parsed.get)
        parsed[largest] = max(0, parsed[largest] + diff)
    return [f"{label}: {parsed[label]}%" for label in wanted]


def clean_company_name(text: str) -> str:
    text = str(text or "").strip(" •-–—\n\t")
    text = re.sub(r"\s*\([^)]*\)\s*", " ", text).strip()
    remove_phrases = [
        "liever ook niet", "ook niet", "niet benaderen", "niet sourcen", "no go",
        "no-go", "eerst checken", "eerst check", "samenwerkingscontracten", "samenwerkingscontract",
        "concurrent", "concurrenten", "die wil hij", "die wil zij", "die wil men",
    ]
    lowered = text.lower()
    for phrase in remove_phrases:
        lowered = lowered.replace(phrase, "")
    # behoud hoofdletters zo veel mogelijk door dezelfde woorden uit originele tekst grof te verwijderen
    text = re.sub(r"(?i)liever ook niet|ook niet|niet benaderen|niet sourcen|no[- ]?go|eerst checken|eerst check|samenwerkingscontracten?|concurrenten?|die wil hij.*|die wil zij.*|die wil men.*", "", text).strip(" :;,-")
    # Pak alleen het eerste deel als er uitleg achter staat.
    text = re.split(r"\s+-\s+|\s+:\s+", text)[0].strip()
    # Veelgemaakte notatie normaliseren.
    replacements = {
        "witteveen en bos": "Witteveen+Bos",
        "witteveen+bos": "Witteveen+Bos",
        "haskoning": "Royal HaskoningDHV",
        "royal haskoning": "Royal HaskoningDHV",
        "antea": "Antea Group",
        "sweco": "Sweco",
    }
    key = text.lower().replace("&", "en").strip()
    return replacements.get(key, text)


def extract_no_go_companies_from_intake(intake_text: str) -> List[str]:
    """Haalt expliciet genoemde no-go/check-eerst organisaties uit intakeblokken."""
    if not intake_text:
        return []
    lines = [ln.strip() for ln in intake_text.splitlines()]
    triggers = ["no-go", "no go", "niet sourcen", "niet benaderen", "samenwerkingscontract", "eerst checken", "liever ook niet"]
    companies: List[str] = []
    collect_next = False
    remaining = 0
    for raw in lines:
        line = raw.strip(" •\t")
        low = line.lower()
        if not line:
            if collect_next:
                remaining -= 1
                if remaining <= 0:
                    collect_next = False
            continue
        has_trigger = any(t in low for t in triggers)
        if has_trigger:
            # Bedrijven kunnen op dezelfde regel of op de regels erna staan.
            collect_next = True
            remaining = 8
            cleaned = clean_company_name(line)
            if cleaned and len(cleaned.split()) <= 5 and cleaned.lower() not in {"samenwerkingscontract", "samenwerkingscontracten"}:
                companies.append(cleaned)
            continue
        if collect_next:
            # Stop bij duidelijke nieuwe vraag/sectie zonder bedrijfsnaam.
            if line.startswith("·") or line.lower().startswith(("meer voorbeeld", "mailtje", "laura", "koen")):
                collect_next = False
                continue
            # Split meerdere bedrijven op komma's of slashes.
            parts = re.split(r",|/|;", line)
            for part in parts:
                cleaned = clean_company_name(part)
                if cleaned and len(cleaned.split()) <= 5:
                    companies.append(cleaned)
            remaining -= 1
            if remaining <= 0:
                collect_next = False
    # Deduplicate, behoud volgorde
    result = []
    for c in companies:
        if c and c not in result:
            result.append(c)
    return result




def strip_bullet_markers(text: str) -> str:
    """Maakt van eventuele bullet-output weer één lopende tekst."""
    text = str(text or "").strip()
    # Verwijder bullets aan het begin van regels en maak er lopende zinnen van.
    lines = [re.sub(r"^\s*[•\-*]\s*", "", ln).strip() for ln in text.splitlines() if ln.strip()]
    if len(lines) > 1:
        text = " ".join(lines)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def limit_words(text: str, max_words: int = 28) -> str:
    words = str(text or "").split()
    if len(words) <= max_words:
        return str(text or "").strip()
    return " ".join(words[:max_words]).rstrip(" ,;")


def presentation_bullets(items: List[str], max_items: int = 3) -> List[str]:
    """Presentation layer: exact maximaal 3 bullets, één onderwerp per bullet."""
    out: List[str] = []
    for item in split_one_topic_per_bullet(items):
        item = re.sub(r"\s+", " ", str(item).strip(" •-\n\t"))
        if not item:
            continue
        item = limit_words(item, 22)
        if item not in out:
            out.append(item)
        if len(out) >= max_items:
            break
    return out


def compact_candidate_bullet(text: str, *, max_words: int = 11) -> str:
    """Maak eisen/voorkeuren compact genoeg voor de kandidaat-slide, zonder betekenis weg te gooien."""
    text = re.sub(r"\s+", " ", str(text or "").strip(" •-\n\t.,;:"))
    if not text:
        return ""
    # Verwijder opvulwoorden die vaak veel ruimte innemen maar weinig toevoegen.
    replacements = [
        (r"\baantoonbare\b", ""),
        (r"\bgedegen\b", ""),
        (r"\bminimaal\s+", ""),
        (r"\bten minste\s+", ""),
        (r"\bervaring met\b", "Ervaring met"),
        (r"\bkennis van\b", "Kennis van"),
        (r"\bin een industriële context\b", "industrieel"),
        (r"\bwerk- en denkniveau\b", "niveau"),
        (r"\bpraktijkervaring\b", "ervaring"),
    ]
    for pattern, repl in replacements:
        text = re.sub(pattern, repl, text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text).strip(" ,;:-")
    text = limit_words(text, max_words)
    if text and not text.endswith((".", ")")):
        text += "."
    return text


def compact_candidate_bullets(items: List[str], max_items: int = 3) -> List[str]:
    out: List[str] = []
    for item in presentation_bullets(items, max_items=max_items):
        compact = compact_candidate_bullet(item)
        if compact and compact not in out:
            out.append(compact)
        if len(out) >= max_items:
            break
    return out


def presentation_summary(text: str) -> str:
    """Intake-samenvatting hoort één mooie lopende tekst te zijn, geen bullets en compact genoeg voor de slide."""
    text = strip_bullet_markers(text)
    text = re.sub(r"\s+", " ", text).strip()
    words = text.split()
    # De intake-dia heeft beperkte ruimte. Houd de tekst concreet, maar maximaal circa 115 woorden.
    if len(words) > 115:
        text = " ".join(words[:115]).rstrip(" ,;") + "."
    return text

def apply_business_rules(data: Dict[str, Any], intake_text: str, linkedin_size: str, vacature_text: str = "", extra_notes: str = "") -> Dict[str, Any]:
    """Harde sprint-0.4 regels die altijd gelden, ook na AI-generatie."""
    data = ensure_core_keys(data)
    b = data.setdefault("basisgegevens", {})
    b["datum"] = date.today().strftime("%d-%m-%Y")

    # Vul klantnaam, vacaturenaam en salaris deterministisch aan als AI ze mist.
    fallback = extract_basis_fallback(vacature_text, intake_text)
    for field in ["klantnaam", "vacaturenaam", "salaris"]:
        if is_empty_or_placeholder(b.get(field, "")) and fallback.get(field):
            b[field] = fallback[field]

    # Salaris staat in de PowerPoint altijd als alleen getal/range, zonder valuta- of toelichtingstekst.
    b["salaris"] = normalize_salary_display(b.get("salaris", ""))

    if linkedin_size.strip():
        data.setdefault("doelgroepanalyse", {})["verwachte_doelgroepgrootte"] = linkedin_size.strip()

    # No-go sourcing: alleen bedrijven, volledig uit intake/extra/AI-output, aangevuld met deterministische extractie.
    k = data.setdefault("kandidaatprofiel", {})
    ai_no_go = clean_list(k.get("no_go_sourcing", []))
    extracted = extract_no_go_companies_from_intake(intake_text)
    merged = []
    for item in ai_no_go + extracted:
        item = clean_company_name(item)
        if item and len(item.split()) <= 5 and item not in merged:
            merged.append(item)
    k["no_go_sourcing"] = merged

    # Presentation layer: intake-samenvatting is één lopende tekst.
    data["intake_samenvatting"] = presentation_summary(data.get("intake_samenvatting", ""))

    # Presentation layer: exact maximaal 3 bullets waar de PowerPoint dit vraagt.
    f = data.setdefault("functieprofiel", {})
    f["taken_verantwoordelijkheden"] = presentation_bullets(f.get("taken_verantwoordelijkheden", []), 3)
    f["usp_functie"] = presentation_bullets(f.get("usp_functie", []), 3)

    k["eisen"] = compact_candidate_bullets(k.get("eisen", []), 3)
    k["voorkeuren"] = compact_candidate_bullets(k.get("voorkeuren", []), 3)

    # Pullfactoren en voorwaarden: één onderwerp per bullet, maximaal 3.
    d = data.setdefault("doelgroepanalyse", {})
    d["pullfactoren"] = presentation_bullets(d.get("pullfactoren", []), 3)
    d["leeftijdsverdeling"] = normalize_age_distribution(d.get("leeftijdsverdeling", []))

    v = data.setdefault("voorwaarden", {})
    v["belangrijkste_arbeidsvoorwaarden"] = presentation_bullets(normalize_conditions(v.get("belangrijkste_arbeidsvoorwaarden", [])), 3)

    # Concurrenten mogen nooit placeholders zijn. Altijd bedrijfsnamen.
    c = data.setdefault("concurrentenanalyse", {})
    current_companies = [clean_company_name(x) for x in clean_list(c.get("bedrijven", []))]
    current_companies = [x for x in current_companies if x and not is_placeholder_company(x)]
    if not current_companies:
        current_companies = infer_competitors_offline(vacature_text, intake_text, current_companies)
    c["bedrijven"] = current_companies
    c["relevant"] = True
    return data


def bullets(items: List[str]) -> str:
    cleaned = clean_list(items)
    return "\n".join(f"• {item}" for item in cleaned)


def plain_lines(items: List[str]) -> str:
    return "\n".join(clean_list(items))


def get_nested(data: Dict[str, Any], path: str, default: Any = "") -> Any:
    cur: Any = data
    for part in path.split("."):
        if isinstance(cur, dict) and part in cur:
            cur = cur[part]
        else:
            return default
    return cur


def replace_text_in_shape(shape, replacements: Dict[str, str]) -> None:
    # Recurse into grouped shapes.
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            replace_text_in_shape(subshape, replacements)
        return

    # Replace in normal text frames.
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for key, value in replacements.items():
                    text = text.replace(key, value)
                run.text = text

    # Replace in tables.
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        for key, value in replacements.items():
                            text = text.replace(key, value)
                        run.text = text


def autofit_shape_text(shape) -> None:
    """Houdt gegenereerde PowerPoints netjes: lange tekst wordt iets kleiner gezet."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            autofit_shape_text(subshape)
        return
    if not (hasattr(shape, "text_frame") and shape.has_text_frame):
        return
    text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
    if not text or "{{" in text:
        return
    length = len(text)
    lines = max(1, text.count("\n") + 1)
    # Alleen lange contentvakken aanpassen; korte titels blijven ongemoeid.
    if length < 120 and lines <= 3:
        return
    if length > 900 or lines > 9:
        size = 8
    elif length > 650 or lines > 7:
        size = 9
    elif length > 420 or lines > 5:
        size = 10
    else:
        size = 11
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def clear_unreplaced_placeholders(shape) -> None:
    """Voorkomt dat {{placeholder}} zichtbaar blijft wanneer data ontbreekt."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            clear_unreplaced_placeholders(subshape)
        return
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = re.sub(r"\{\{[^}]+\}\}", "", run.text)
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = re.sub(r"\{\{[^}]+\}\}", "", run.text)


def generate_pptx(data: Dict[str, Any]) -> bytes:
    template_path = get_template_path()
    prs = Presentation(str(template_path))
    delete_slides_by_exact_title(prs, {"AANPAK"})
    afspraken = data.get("afspraken") or []
    concurrenten = get_nested(data, "concurrentenanalyse.bedrijven", [])
    concurrenten_text = bullets(concurrenten) or get_nested(data, "concurrentenanalyse.toelichting", "")

    replacements = {
        "{{klantnaam}}": get_nested(data, "basisgegevens.klantnaam"),
        "{{vacaturenaam}}": get_nested(data, "basisgegevens.vacaturenaam"),
        "{{datum}}": get_nested(data, "basisgegevens.datum") or date.today().strftime("%d-%m-%Y"),
        "{{intake_samenvatting}}": data.get("intake_samenvatting", ""),
        "{{sourcingplan_strategie}}": get_nested(data, "sourcingplan.strategie"),
        "{{sourcingplan_doelgroep}}": get_nested(data, "sourcingplan.doelgroep"),
        "{{concurrentenanalyse}}": concurrenten_text,
        "{{zoekrichting}}": bullets(get_nested(data, "sourcingplan.zoekrichting", [])),
        "{{aanpak_toelichting}}": get_nested(data, "sourcingplan.toelichting"),
        "{{doelgroep_titel}}": get_nested(data, "doelgroepanalyse.doelgroep_titel") or get_nested(data, "basisgegevens.vacaturenaam"),
        "{{taken_verantwoordelijkheden}}": bullets(get_nested(data, "functieprofiel.taken_verantwoordelijkheden", [])),
        "{{eisen}}": bullets(get_nested(data, "kandidaatprofiel.eisen", [])),
        "{{voorkeuren}}": bullets(get_nested(data, "kandidaatprofiel.voorkeuren", [])),
        "{{no_go_sourcing}}": bullets(get_nested(data, "kandidaatprofiel.no_go_sourcing", [])),
        "{{doelgroepgrootte}}": get_nested(data, "doelgroepanalyse.verwachte_doelgroepgrootte"),
        "{{doelgroep_regio}}": get_nested(data, "doelgroepanalyse.regio") or "Nederland",
        "{{salaris}}": get_nested(data, "basisgegevens.salaris"),
        "{{locatie}}": get_nested(data, "basisgegevens.locatie"),
        "{{uren}}": get_nested(data, "basisgegevens.uren"),
        "{{usp_functie}}": bullets(get_nested(data, "functieprofiel.usp_functie", [])),
        "{{pullfactoren}}": bullets(get_nested(data, "doelgroepanalyse.pullfactoren", [])),
        "{{belangrijkste_arbeidsvoorwaarden}}": bullets(get_nested(data, "voorwaarden.belangrijkste_arbeidsvoorwaarden", [])),
        "{{geslacht_man}}": get_nested(data, "doelgroepanalyse.geslacht.man"),
        "{{geslacht_vrouw}}": get_nested(data, "doelgroepanalyse.geslacht.vrouw"),
        "{{leeftijdsverdeling}}": bullets(get_nested(data, "doelgroepanalyse.leeftijdsverdeling", [])),
        "{{afspraken_1}}": afspraken[0] if len(afspraken) > 0 else "",
        "{{afspraken_2}}": afspraken[1] if len(afspraken) > 1 else "",
        "{{afspraken_3}}": afspraken[2] if len(afspraken) > 2 else "",
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
            clear_unreplaced_placeholders(shape)
            autofit_shape_text(shape)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return Path(tmp.name).read_bytes()



# -----------------------------------------------------------------------------
# v1.2 Presentation Engine
# -----------------------------------------------------------------------------

FONT_NAME = "Poppins"
DARK_BLUE = RGBColor(0, 32, 96)
BODY_COLOR = RGBColor(0, 0, 0)
ACCENT_BLUE = "#005B99"

TEXT_STYLES = {
    # Cooble-template stijl. Groottes zijn in pt en afgestemd op de originele template.
    "cover_title": {"size": 104, "bold": True, "color": DARK_BLUE},
    "cover_footer": {"size": 28, "bold": True, "color": DARK_BLUE},
    "slide_title": {"size": 54, "bold": True, "color": DARK_BLUE},
    "big_title": {"size": 76, "bold": True, "color": DARK_BLUE},
    "date": {"size": 26, "bold": False, "color": BODY_COLOR},
    "customer": {"size": 24, "bold": False, "color": BODY_COLOR},
    "intake": {"size": 21, "bold": False, "color": BODY_COLOR},
    "body": {"size": 17, "bold": False, "color": BODY_COLOR},
    "body_small": {"size": 15, "bold": False, "color": BODY_COLOR},
    "candidate_bullet": {"size": 14, "bold": False, "color": BODY_COLOR},
    "heading": {"size": 34, "bold": True, "color": DARK_BLUE},
    "subtitle": {"size": 26, "bold": True, "color": DARK_BLUE},
    "section_heading": {"size": 19, "bold": True, "color": DARK_BLUE},
    "small_heading": {"size": 17, "bold": True, "color": DARK_BLUE},
    "metric": {"size": 88, "bold": True, "color": BODY_COLOR},
    "percentage": {"size": 24, "bold": True, "color": BODY_COLOR},
}

BULLET_PLACEHOLDERS = {
    "{{taken_verantwoordelijkheden}}": ("functieprofiel.taken_verantwoordelijkheden", "body_small"),
    "{{eisen}}": ("kandidaatprofiel.eisen", "candidate_bullet"),
    "{{voorkeuren}}": ("kandidaatprofiel.voorkeuren", "candidate_bullet"),
    "{{usp_functie}}": ("functieprofiel.usp_functie", "body"),
    "{{no_go_sourcing}}": ("kandidaatprofiel.no_go_sourcing", "body"),
    "{{pullfactoren}}": ("doelgroepanalyse.pullfactoren", "body"),
    "{{belangrijkste_arbeidsvoorwaarden}}": ("voorwaarden.belangrijkste_arbeidsvoorwaarden", "body"),
    "{{concurrentenanalyse}}": ("concurrentenanalyse.bedrijven", "body"),
    "{{zoekrichting}}": ("sourcingplan.zoekrichting", "body"),
}

PLAIN_PLACEHOLDERS = {
    "{{klantnaam}}": ("basisgegevens.klantnaam", "customer"),
    "{{datum}}": ("basisgegevens.datum", "date"),
    "{{intake_samenvatting}}": ("intake_samenvatting", "intake"),
    "{{sourcingplan_strategie}}": ("sourcingplan.strategie", "body"),
    "{{sourcingplan_doelgroep}}": ("sourcingplan.doelgroep", "body"),
    "{{aanpak_toelichting}}": ("sourcingplan.toelichting", "body"),
    "{{doelgroepgrootte}}": ("doelgroepanalyse.verwachte_doelgroepgrootte", "metric"),
    "{{geslacht_man}}": ("doelgroepanalyse.geslacht.man", "percentage"),
    "{{geslacht_vrouw}}": ("doelgroepanalyse.geslacht.vrouw", "percentage"),
    "{{afspraken_1}}": ("afspraken.0", "body"),
    "{{afspraken_2}}": ("afspraken.1", "body"),
    "{{afspraken_3}}": ("afspraken.2", "body"),
}


def get_nested_v12(data: Dict[str, Any], path: str, default: Any = "") -> Any:
    cur: Any = data
    for part in path.split("."):
        if isinstance(cur, list):
            if part.isdigit() and int(part) < len(cur):
                cur = cur[int(part)]
            else:
                return default
        elif isinstance(cur, dict) and part in cur:
            cur = cur[part]
        else:
            return default
    return cur


def full_text(shape) -> str:
    if not (hasattr(shape, "text_frame") and shape.has_text_frame):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def clear_tf(tf) -> None:
    tf.clear()
    tf.word_wrap = True
    # marges iets rustiger zoals PowerPoint-template
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def font_size_for_text(style_name: str, text: str, lines: int = 1) -> int:
    base = TEXT_STYLES.get(style_name, TEXT_STYLES["body"])["size"]
    words = len(str(text).split())
    chars = len(str(text))
    if style_name == "intake":
        # Intake moet compact blijven zodat titel en subtitel vrij blijven.
        if words > 145 or chars > 920:
            return 16
        if words > 120 or chars > 760:
            return 18
        if words > 95 or chars > 620:
            return 19
        return base
    if style_name == "candidate_bullet":
        if lines >= 3 or chars > 210:
            return 12
        return base
    if style_name in {"body", "body_small"}:
        if lines >= 7 or chars > 520:
            return 14 if style_name == "body_small" else 15
        if lines >= 5 or chars > 380:
            return 15 if style_name == "body_small" else 16
        if chars > 260:
            return 16
    return base


def apply_run_style(run, style_name: str, *, size_override: int | None = None) -> None:
    stl = TEXT_STYLES.get(style_name, TEXT_STYLES["body"])
    run.font.name = FONT_NAME
    run.font.size = Pt(size_override or stl["size"])
    run.font.bold = stl.get("bold", False)
    run.font.color.rgb = stl.get("color", BODY_COLOR)


def set_plain_text(shape, text: str, style_name: str) -> None:
    tf = shape.text_frame
    clear_tf(tf)
    text = str(text or "").strip()
    p = tf.paragraphs[0]
    p.alignment = None
    run = p.add_run()
    run.text = text
    size = font_size_for_text(style_name, text)
    apply_run_style(run, style_name, size_override=size)
    for par in tf.paragraphs:
        par.space_after = Pt(0)
        par.space_before = Pt(0)


def set_bullet_list(shape, items: List[str], style_name: str = "body") -> None:
    tf = shape.text_frame
    clear_tf(tf)
    clean = clean_list(items)[:8]
    if not clean:
        return
    combined = "\n".join(clean)
    size = font_size_for_text(style_name, combined, len(clean))
    for idx, item in enumerate(clean):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        p.level = 0
        p.space_after = Pt(3 if style_name == "candidate_bullet" else (6 if len(clean) <= 3 else 4))
        p.space_before = Pt(0)
        # Gebruik een tekst-bullet zodat het stabiel blijft in PowerPoint en template-stijl benadert.
        run = p.add_run()
        run.text = str(item).strip()
        apply_run_style(run, style_name, size_override=size)
        # echte bullet-properties via XML worden soms genegeerd door PowerPoint, daarom zetten we een bullet prefix veilig.
        p._p.get_or_add_pPr()
        if not run.text.startswith("•"):
            run.text = "• " + run.text


def set_mixed_text_frame(shape, replacements: Dict[str, str], style_name: str = "body") -> None:
    # Voor tekstvakken met labels + placeholders, zoals salaris/locatie/uren.
    original = full_text(shape)
    text = original
    for key, value in replacements.items():
        text = text.replace(key, str(value or ""))
    tf = shape.text_frame
    clear_tf(tf)
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    if not lines:
        lines = [text.strip()]
    size = font_size_for_text(style_name, "\n".join(lines), len(lines))
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        apply_run_style(run, style_name, size_override=size)
        p.space_after = Pt(4)


def parse_age_items(items: List[str]) -> tuple[list[str], list[float]]:
    labels, values = [], []
    for raw in clean_list(items):
        text = str(raw)
        m = re.search(r"(\d{1,3})\s*%", text)
        value = float(m.group(1)) if m else 0.0
        label = re.sub(r"[:\-–]?\s*\d{1,3}\s*%", "", text).strip(" :-–")
        if not label:
            label = text.strip()
        labels.append(label)
        values.append(value)
    if not labels or sum(values) == 0:
        labels = ["25-34", "35-44", "45-54", "55+"]
        values = [30, 35, 25, 10]
    return labels[:5], values[:5]


def create_age_chart_image(items: List[str]) -> BytesIO:
    """Maakt leeftijdsverdeling als horizontale balken, zoals de aangeleverde referentie."""
    import matplotlib.pyplot as plt
    labels, values = parse_age_items(items)
    # Standaardiseer de categorieën iets compacter voor de slide.
    labels = [str(x).replace(" ", "") for x in labels]
    max_value = max(max(values), 100)

    fig, ax = plt.subplots(figsize=(5.8, 3.0), dpi=180)
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    ax.set_xlim(-28, 118)
    ax.set_ylim(-0.6, len(labels) - 0.4)
    ax.axis("off")

    accent = "#5A4BD8"  # paarse balk zoals de aangeleverde referentie
    track = "#EDEFF2"
    label_color = "#111111"

    for i, (label, value) in enumerate(zip(labels, values)):
        y = len(labels) - 1 - i
        ax.text(-26, y, label, va="center", ha="left", fontsize=10.5, color=label_color)
        ax.plot([0, 100], [y, y], color=track, linewidth=15, solid_capstyle="round")
        if value > 0:
            ax.plot([0, min(value, 100)], [y, y], color=accent, linewidth=15, solid_capstyle="round")
        ax.text(113, y, f"{int(round(value))}%", va="center", ha="right", fontsize=11.5, fontweight="bold", color=label_color)

    buf = BytesIO()
    plt.savefig(buf, format="png", transparent=True, bbox_inches="tight", pad_inches=0.03)
    plt.close(fig)
    buf.seek(0)
    return buf

def render_shape_v12(slide, shape, data: Dict[str, Any], replacements: Dict[str, str]) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            render_shape_v12(slide, subshape, data, replacements)
        return
    if not (hasattr(shape, "text_frame") and shape.has_text_frame):
        return
    original = full_text(shape)
    if not original:
        return

    if "{{leeftijdsverdeling}}" in original:
        # Tekstplaceholder leegmaken en op dezelfde dia een echte afbeelding/infographic plaatsen.
        shape.text_frame.clear()
        img = create_age_chart_image(get_nested_v12(data, "doelgroepanalyse.leeftijdsverdeling", []))
        # Vaste chart-zone rechts op de leeftijdsdia, gecentreerd onder de kop.
        slide.shapes.add_picture(img, Emu(9300000), Emu(5100000), width=Emu(6900000), height=Emu(3000000))
        return

    stripped = original.strip()

    # Vaste templatekoppen opnieuw stylen volgens Cooble-template.
    upper = stripped.upper().replace("\n", " ").strip()
    if upper == "START DOCUMENT":
        set_plain_text(shape, "START\nDOCUMENT", "cover_title")
        return
    if upper == "INTAKE":
        set_plain_text(shape, "INTAKE", "big_title")
        return
    if upper in {"TAKEN & VERANTWOORDELIJKHEDEN", "EISEN", "VOORKEUREN", "USP'S VAN DE FUNCTIE", "NO GO SOURCING", "BELANGRIJKSTE PULLFACTOREN", "BELANGRIJKSTE ARBEIDSVOORWAARDEN", "GESLACHT", "LEEFTIJDSVERDELING (IN JAREN)"}:
        set_plain_text(shape, stripped, "small_heading" if len(stripped) > 24 else "section_heading")
        return
    if upper in {"KANDIDAAT", "AFSPRAKEN", "HET PROCES VAN COOBLE", "DOELGROEP ANALYSE"}:
        set_plain_text(shape, stripped, "slide_title")
        return

    # Op de eerste slide staat {{klantnaam}} onderin; daar wil je de vacaturetitel zien.
    if stripped == "{{klantnaam}}" and getattr(shape, "top", 0) > Emu(8000000):
        set_plain_text(shape, str(get_nested_v12(data, "basisgegevens.vacaturenaam", "")), "cover_footer")
        return

    # Kleine subtitel onder INTAKE; voorkom overlap met de samenvatting.
    if stripped == "{{vacaturenaam}}" and getattr(shape, "height", 0) < Emu(900000):
        set_plain_text(shape, str(get_nested_v12(data, "basisgegevens.vacaturenaam", "")), "subtitle")
        return

    for placeholder, (path, style) in BULLET_PLACEHOLDERS.items():
        if placeholder in stripped and stripped == placeholder:
            set_bullet_list(shape, get_nested_v12(data, path, []), style)
            return

    for placeholder, (path, style) in PLAIN_PLACEHOLDERS.items():
        if placeholder in stripped and stripped == placeholder:
            set_plain_text(shape, get_nested_v12(data, path, ""), style)
            return

    # Speciale titel: geen "DOELGROEP:" meer, alleen de concrete vacature/doelgroep.
    if "{{doelgroep_titel}}" in original:
        text = str(get_nested_v12(data, "doelgroepanalyse.doelgroep_titel") or get_nested_v12(data, "basisgegevens.vacaturenaam", ""))
        set_plain_text(shape, text, "slide_title")
        return

    if "{{vacaturenaam}}" in original:
        text = original.replace("{{vacaturenaam}}", str(get_nested_v12(data, "basisgegevens.vacaturenaam", "")))
        style = "subtitle" if len(text) < 55 else "section_heading"
        set_plain_text(shape, text, style)
        return

    if "{{" in original:
        set_mixed_text_frame(shape, replacements, "body_small")
        return

    # Vaste koppen worden bovenin al gestyled.



def delete_slides_by_exact_title(prs, titles: set[str]) -> None:
    """Verwijdert slides uit de Cooble-template, bijvoorbeeld de oude AANPAK-slide."""
    slides_to_delete = []
    for idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip().upper())
        if any(t in titles for t in texts):
            slides_to_delete.append(idx)
    # python-pptx heeft geen publieke delete API; dit is de gebruikelijke XML-route.
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    for idx in sorted(slides_to_delete, reverse=True):
        r_id = sld_id_lst[idx].rId
        prs.part.drop_rel(r_id)
        del sld_id_lst[idx]

def generate_pptx(data: Dict[str, Any]) -> bytes:
    """v1.3: Cooble-template behouden, AANPAK-slide verwijderd en betere PowerPoint-rendering."""
    template_path = get_template_path()
    prs = Presentation(str(template_path))
    delete_slides_by_exact_title(prs, {"AANPAK"})
    afspraken = data.get("afspraken") or []
    concurrenten = get_nested_v12(data, "concurrentenanalyse.bedrijven", [])
    concurrenten_text = bullets(concurrenten) or get_nested_v12(data, "concurrentenanalyse.toelichting", "")

    replacements = {
        "{{klantnaam}}": get_nested_v12(data, "basisgegevens.klantnaam"),
        "{{vacaturenaam}}": get_nested_v12(data, "basisgegevens.vacaturenaam"),
        "{{datum}}": get_nested_v12(data, "basisgegevens.datum") or date.today().strftime("%d-%m-%Y"),
        "{{intake_samenvatting}}": data.get("intake_samenvatting", ""),
        "{{sourcingplan_strategie}}": get_nested_v12(data, "sourcingplan.strategie"),
        "{{sourcingplan_doelgroep}}": get_nested_v12(data, "sourcingplan.doelgroep"),
        "{{concurrentenanalyse}}": concurrenten_text,
        "{{zoekrichting}}": bullets(get_nested_v12(data, "sourcingplan.zoekrichting", [])),
        "{{aanpak_toelichting}}": get_nested_v12(data, "sourcingplan.toelichting"),
        "{{doelgroep_titel}}": get_nested_v12(data, "doelgroepanalyse.doelgroep_titel") or get_nested_v12(data, "basisgegevens.vacaturenaam"),
        "{{taken_verantwoordelijkheden}}": bullets(get_nested_v12(data, "functieprofiel.taken_verantwoordelijkheden", [])),
        "{{eisen}}": bullets(get_nested_v12(data, "kandidaatprofiel.eisen", [])),
        "{{voorkeuren}}": bullets(get_nested_v12(data, "kandidaatprofiel.voorkeuren", [])),
        "{{no_go_sourcing}}": bullets(get_nested_v12(data, "kandidaatprofiel.no_go_sourcing", [])),
        "{{doelgroepgrootte}}": get_nested_v12(data, "doelgroepanalyse.verwachte_doelgroepgrootte"),
        "{{doelgroep_regio}}": get_nested_v12(data, "doelgroepanalyse.regio") or "Nederland",
        "{{salaris}}": get_nested_v12(data, "basisgegevens.salaris"),
        "{{locatie}}": get_nested_v12(data, "basisgegevens.locatie"),
        "{{uren}}": get_nested_v12(data, "basisgegevens.uren"),
        "{{usp_functie}}": bullets(get_nested_v12(data, "functieprofiel.usp_functie", [])),
        "{{pullfactoren}}": bullets(get_nested_v12(data, "doelgroepanalyse.pullfactoren", [])),
        "{{belangrijkste_arbeidsvoorwaarden}}": bullets(get_nested_v12(data, "voorwaarden.belangrijkste_arbeidsvoorwaarden", [])),
        "{{geslacht_man}}": get_nested_v12(data, "doelgroepanalyse.geslacht.man"),
        "{{geslacht_vrouw}}": get_nested_v12(data, "doelgroepanalyse.geslacht.vrouw"),
        "{{leeftijdsverdeling}}": bullets(get_nested_v12(data, "doelgroepanalyse.leeftijdsverdeling", [])),
        "{{afspraken_1}}": afspraken[0] if len(afspraken) > 0 else "",
        "{{afspraken_2}}": afspraken[1] if len(afspraken) > 1 else "",
        "{{afspraken_3}}": afspraken[2] if len(afspraken) > 2 else "",
    }

    for slide in prs.slides:
        for shape in list(slide.shapes):
            render_shape_v12(slide, shape, data, replacements)
        for shape in slide.shapes:
            clear_unreplaced_placeholders(shape)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return Path(tmp.name).read_bytes()

def schema_hint() -> str:
    return """
Geef uitsluitend geldige JSON terug in exact deze structuur:
{
  "basisgegevens": {
    "klantnaam": "",
    "vacaturenaam": "",
    "datum": "",
    "locatie": "",
    "uren": "",
    "salaris": ""
  },
  "intake_samenvatting": "",
  "functieprofiel": {
    "taken_verantwoordelijkheden": ["", "", ""],
    "usp_functie": ["", "", ""]
  },
  "kandidaatprofiel": {
    "eisen": ["", "", ""],
    "voorkeuren": ["", "", ""],
    "no_go_sourcing": []
  },
  "voorwaarden": {
    "belangrijkste_arbeidsvoorwaarden": ["", "", ""]
  },
  "doelgroepanalyse": {
    "doelgroep_titel": "",
    "verwachte_doelgroepgrootte": "",
    "regio": "Nederland",
    "pullfactoren": ["", "", ""],
    "geslacht": {"man": "", "vrouw": ""},
    "leeftijdsverdeling": ["", "", "", ""]
  },
  "sourcingplan": {
    "doelgroep": "",
    "strategie": "",
    "belangrijkste_functietitels": [],
    "zoekrichting": [],
    "toelichting": ""
  },
  "concurrentenanalyse": {
    "relevant": true,
    "bedrijven": [],
    "toelichting": ""
  },
  "afspraken": [],
  "kwaliteitscontrole": {
    "ontbrekende_informatie": [],
    "aannames": [],
    "waarschuwingen": []
  }
}
""".strip()


def build_prompt(vacature: str, intake: str, linkedin_size: str, extra: str) -> str:
    return f"""
Je bent een senior recruitment consultant en arbeidsmarktanalist. Maak compacte inhoud voor een Cooble startdocument.

Belangrijke regels:
- Output is Nederlands.
- Eén intake = één vacature.
- Formuleer kort en bondig, PowerPoint-stijl.
- Concurrentenanalyse is altijd relevant en altijd op bedrijfsniveau. Gebruik echte bedrijfsnamen. Nooit placeholders zoals Bedrijf A, Bedrijf B of Concurrent 1.
- Vrij formuleren mag, maar niet fantaseren.
- Datum: laat leeg of gebruik vandaag; de app overschrijft dit altijd met de generatiedatum.
- No-go sourcing: geef uitsluitend bedrijfsnamen terug. Neem álle no-go/check-eerst organisaties uit de intake over. Geen toelichting, geen zinnen. Laat geen enkel bedrijf weg.
- Pullfactoren zijn extern: bepaal ze vanuit arbeidsmarkt/doelgroep en internetonderzoek, niet uit de vacaturetekst.
- Belangrijkste arbeidsvoorwaarden: niet uit de vacaturetekst halen. Onderzoek extern/arbeidsmarktgericht welke arbeidsvoorwaarden de doelgroep belangrijk vindt. Gebruik generieke labels zoals "Vakantiedagen", niet "29 vakantiedagen".
- Eén bullet = één onderwerp. Lijsten voor taken/eisen/voorkeuren/USP/pullfactoren/arbeidsvoorwaarden bevatten precies 3 items.
- Eisen en voorkeuren zijn compact: maximaal 10 woorden per bullet, zonder opvulwoorden zoals 'relevante ervaring'. Combineer nooit meerdere onderwerpen in één bullet.
- Intake is leidend boven vacaturetekst.
- Extra opmerkingen zijn leidend boven alles.
- Houd lijsten kort: meestal precies 3 bullets, behalve no-go sourcing en concurrenten; die mogen meer bedrijven bevatten.
- Gebruik zelfverzekerde labels zoals "Hybride werken", niet "waarschijnlijk hybride werken".
- Leeftijdsverdeling: geef categorie én percentage, bijvoorbeeld "25-34: 30%".
- Als doelgroepgrootte uit LinkedIn is ingevuld, gebruik die waarde letterlijk.
- Klantnaam en vacaturenaam moeten altijd gevuld zijn. Haal klantnaam uit intake/vacaturetekst. Haal vacaturenaam uit intake/vacaturetitel.
- Salaris: als er een schaal of salarisrange in intake/vacature staat, neem die concreet over. Gebruik alleen een getal/range; bij een salarisschaal mag "Schaal" ervoor staan. Gebruik niet "in overleg" als er schalen of bedragen staan.
- Intake_samenvatting: schrijf concreet maar compact. Dit veld moet in één dia duidelijk maken waar we naar zoeken, inclusief aanleiding, focus, nuances, wat juist niet past en nadruk uit de intake. Richtlijn: 70-100 woorden.
- Taken & verantwoordelijkheden: vermijd generieke bullets. Benoem de inhoudelijke context, doelgroep/klanttype, projecten of domein.
- Eisen: vermijd "relevante ervaring". Schrijf ervaring waarmee, bijvoorbeeld "ervaring met industriële waterprojecten".
- Doelgroep: zeer concreet. Gebruik vacaturetitel, domein, senioriteit, sector en relevante achtergrond. Nooit generiek zoals "kandidaten met relevante advieservaring".
- Concurrenten: doe internetonderzoek en geef echte bedrijven waar deze doelgroep werkt of vandaan kan komen.

{schema_hint()}

VACATURETEKST:
{vacature[:30000]}

INTAKE NOTES:
{intake[:30000]}

DOELGROEPGROOTTE GEVONDEN OP LINKEDIN:
{linkedin_size[:500]}

EXTRA OPMERKINGEN:
{extra[:10000]}
"""


def extract_json(text: str) -> Dict[str, Any]:
    text = (text or "").strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?", "", text).strip()
        text = re.sub(r"```$", "", text).strip()
    try:
        return json.loads(text)
    except Exception:
        match = re.search(r"\{.*\}", text, flags=re.S)
        if match:
            return json.loads(match.group(0))
        raise



def call_openai_json(prompt: str, *, use_web: bool = False, system: str = "Je geeft uitsluitend geldige JSON terug. Geen markdown.") -> Dict[str, Any]:
    """Centrale OpenAI-call. Probeert Responses API met optioneel webonderzoek, valt terug op Chat Completions."""
    api_key = st.secrets.get("OPENAI_API_KEY", "")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY ontbreekt in Streamlit Secrets.")
    client = OpenAI(api_key=api_key)
    model = st.secrets.get("OPENAI_MODEL", DEFAULT_MODEL)

    if use_web:
        try:
            response = client.responses.create(
                model=model,
                input=prompt,
                tools=[{"type": "web_search_preview"}],
            )
            return extract_json(response.output_text)
        except Exception as first_error:
            # Niet stoppen: sommige modellen/accounts ondersteunen web_search_preview niet.
            data = call_openai_json(
                prompt + "\n\nLET OP: web_search_preview was niet beschikbaar. Gebruik algemene arbeidsmarktkennis, maar blijf concreet.",
                use_web=False,
                system=system,
            )
            data.setdefault("_meta", {})["web_search_warning"] = str(first_error)
            return data

    chat = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        response_format={"type": "json_object"},
    )
    return extract_json(chat.choices[0].message.content or "{}")


def build_fact_extraction_prompt(vacature: str, intake: str, extra: str) -> str:
    return f"""
Je bent een nauwkeurige recruitment-analist. Haal uitsluitend FEITEN uit de vacaturetekst, intake en extra opmerkingen.
Niet interpreteren, niet mooier maken, niet aanvullen.

Regels:
- Intake is leidend boven vacaturetekst.
- Extra opmerkingen zijn leidend boven alles.
- Neem salaris/schalen concreet over als ze genoemd worden.
- Neem alle no-go/check-eerst organisaties uit de intake over als losse bedrijfsnamen.
- Als een veld ontbreekt, gebruik een lege string of lege lijst.

Geef uitsluitend JSON terug:
{{
  "klantnaam": "",
  "vacaturenaam": "",
  "locatie": "",
  "uren": "",
  "salaris": "",
  "aanleiding_vacature": "",
  "manager_nadruk": [],
  "nuances": [],
  "wat_past_niet": [],
  "taken_feiten": [],
  "eisen_feiten": [],
  "voorkeuren_feiten": [],
  "usp_feiten": [],
  "arbeidsvoorwaarden_uit_vacature": [],
  "no_go_bedrijven": [],
  "afspraken": []
}}

VACATURETEKST:
{vacature[:30000]}

INTAKE NOTES:
{intake[:30000]}

EXTRA OPMERKINGEN:
{extra[:10000]}
""".strip()


def build_research_prompt(facts: Dict[str, Any], linkedin_size: str, vacature: str, intake: str) -> str:
    functie = facts.get("vacaturenaam") or first_nonempty_line(vacature)
    klant = facts.get("klantnaam", "")
    locatie = facts.get("locatie", "Nederland")
    nuances = facts.get("nuances", [])
    manager_nadruk = facts.get("manager_nadruk", [])
    return f"""
Je bent een recruitment researcher voor de Nederlandse arbeidsmarkt.
Onderzoek de doelgroep voor deze vacature en geef concrete, niet-generieke onderzoeksoutput.

Belangrijke regels:
- Pullfactoren zijn extern: baseer ze op doelgroep/arbeidsmarkt, NIET op de vacaturetekst.
- Arbeidsvoorwaarden zijn extern: doe internetonderzoek naar wat deze doelgroep belangrijk vindt. Gebruik de vacaturetekst of intake NIET als bron voor dit veld, behalve als extra opmerkingen iets expliciet verplichten. Gebruik generieke arbeidsmarktlabels, geen concrete waarden uit de vacature.
- Concurrentenanalyse is altijd relevant en altijd op bedrijfsniveau.
- Geef echte bedrijfsnamen. Nooit Bedrijf A/B/C, Concurrent 1, Organisatie X.
- Doelgroepomschrijving moet specifiek zijn voor functie, domein, senioriteit en sector.
- Eén bullet = één onderwerp. Lijsten voor taken/eisen/voorkeuren/USP/pullfactoren/arbeidsvoorwaarden bevatten precies 3 items.
- Eisen en voorkeuren zijn compact: maximaal 10 woorden per bullet, zonder opvulwoorden zoals 'relevante ervaring'.
- Geen voorzichtige taal zoals "waarschijnlijk" of "mogelijk" in de uiteindelijke labels.
- Als doelgroepgrootte uit LinkedIn is ingevuld, neem die letterlijk over.

Context uit de intake/vacature:
Klant: {klant}
Functie: {functie}
Locatie/regio: {locatie}
Doelgroepgrootte LinkedIn: {linkedin_size}
Nuances: {json.dumps(nuances, ensure_ascii=False)}
Manager nadruk: {json.dumps(manager_nadruk, ensure_ascii=False)}
No-go/check-eerst bedrijven: {json.dumps(facts.get('no_go_bedrijven', []), ensure_ascii=False)}

Geef uitsluitend JSON terug:
{{
  "doelgroep_titel": "",
  "doelgroep_omschrijving": "",
  "verwachte_doelgroepgrootte": "",
  "belangrijkste_functietitels": [],
  "pullfactoren": ["", "", ""],
  "belangrijkste_arbeidsvoorwaarden": ["", "", ""],
  "concurrenten_bedrijven": [],
  "zoekrichting": [],
  "geslacht": {{"man": "", "vrouw": ""}},
  "leeftijdsverdeling": ["25-34: %", "35-44: %", "45-54: %", "55+: %"],
  "research_toelichting": ""
}}
""".strip()


def build_writer_prompt(facts: Dict[str, Any], research: Dict[str, Any], vacature: str, intake: str, linkedin_size: str, extra: str) -> str:
    return f"""
Je bent een senior recruitment consultant van Cooble/Sinvae. Schrijf de definitieve PowerPoint-content voor een startdocument.

Strenge schrijfrichtlijnen:
- Nederlands.
- Kort en concreet, PowerPoint-stijl.
- Vrij formuleren, maar niet fantaseren.
- Intake is leidend boven vacaturetekst.
- Extra opmerkingen zijn leidend boven alles.
- Gebruik de feitenextractie en research als bron; schrijf niet opnieuw generiek vanuit de vacature.
- Intake_samenvatting: 80-110 woorden als één mooie lopende tekst. Geen bullets, geen opsomming. Deze dia moet zelfstandig duidelijk maken waar we naar zoeken, inclusief aanleiding, focus, nuances, nadruk uit intake en wat juist niet past.
- Taken: precies 3 bullets, concreet voor deze rol. Benoem domein, klanttype, projecttype of inhoudelijke context.
- Eisen: precies 3 bullets. Nooit "relevante ervaring". Schrijf ervaring waarmee.
- Doelgroep: specifiek voor deze functie, sector, senioriteit en domein.
- Pullfactoren: extern en arbeidsmarktgericht, niet uit vacaturetekst.
- Arbeidsvoorwaarden: gebruik uitsluitend het aparte online onderzoek naar wat deze doelgroep belangrijk vindt WANNEER ZIJ IN DIENST ZIJN. Niet uit vacaturetekst of intake overnemen. Kies 3 arbeidsvoorwaardencategorieën, bijvoorbeeld Salaris, Pensioenregeling, Vakantiedagen, Hybride werken, Mobiliteit of Ontwikkelmogelijkheden.
- Pullfactoren: gebruik uitsluitend het aparte online onderzoek naar wat deze doelgroep IN BEWEGING BRENGT en wat zij graag terugzien in een vacature. Pullfactoren zijn géén arbeidsvoorwaardenlijst en mogen niet uit vacaturetekst of intake worden afgeleid.
- Salaris: toon uitsluitend het getal of de range, zonder €, bruto, per maand, schaal of andere woorden.
- Concurrenten: echte bedrijfsnamen op bedrijfsniveau.
- No-go sourcing: uitsluitend bedrijfsnamen uit feitenextractie. Niets toevoegen, niets weglaten.
- Eén bullet = één onderwerp. Lijsten voor taken/eisen/voorkeuren/USP/pullfactoren/arbeidsvoorwaarden bevatten precies 3 items.
- Eisen en voorkeuren zijn compact: maximaal 10 woorden per bullet, zonder opvulwoorden zoals 'relevante ervaring'.
- Vermijd generieke termen: relevante ervaring, passende kandidaat, dynamische omgeving, goede communicatieve vaardigheden, inhoudelijk specialistisch domein, spin in het web.

{schema_hint()}

FEITENEXTRACTIE:
{json.dumps(facts, ensure_ascii=False, indent=2)}

ARBEIDSMARKT- EN DOELGROEPRESEARCH:
{json.dumps(research, ensure_ascii=False, indent=2)}

DOELGROEPGROOTTE GEVONDEN OP LINKEDIN:
{linkedin_size[:500]}

EXTRA OPMERKINGEN:
{extra[:10000]}

TER CONTROLE - ORIGINELE VACATURE:
{vacature[:20000]}

TER CONTROLE - ORIGINELE INTAKE:
{intake[:20000]}
""".strip()


GENERIC_PHRASES = [
    "relevante ervaring",
    "passende kandidaat",
    "dynamische omgeving",
    "goede communicatieve vaardigheden",
    "inhoudelijk specialistisch domein",
    "spin in het web",
    "uitdagende functie",
    "veelzijdige functie",
    "marktconform salaris",
]


def collect_generic_issues(data: Dict[str, Any]) -> List[str]:
    issues: List[str] = []
    paths = {
        "intake_samenvatting": data.get("intake_samenvatting", ""),
        "taken": " | ".join(get_nested(data, "functieprofiel.taken_verantwoordelijkheden", [])),
        "eisen": " | ".join(get_nested(data, "kandidaatprofiel.eisen", [])),
        "voorkeuren": " | ".join(get_nested(data, "kandidaatprofiel.voorkeuren", [])),
        "doelgroep": get_nested(data, "sourcingplan.doelgroep", ""),
    }
    for label, text in paths.items():
        low = str(text).lower()
        for phrase in GENERIC_PHRASES:
            if phrase in low:
                issues.append(f"{label}: vermijd '{phrase}'")
    # Check placeholders bij concurrenten
    for item in get_nested(data, "concurrentenanalyse.bedrijven", []):
        if is_placeholder_company(item):
            issues.append("concurrenten: placeholder-bedrijf gevonden")
    return issues


def build_refine_prompt(data: Dict[str, Any], facts: Dict[str, Any], research: Dict[str, Any], issues: List[str]) -> str:
    return f"""
Verbeter deze startdocument-JSON. Los alleen de genoemde kwaliteitsissues op.
Behoud de JSON-structuur exact. Voeg geen nieuwe feiten toe die niet uit feitenextractie of research komen.

Issues:
{json.dumps(issues, ensure_ascii=False, indent=2)}

Regels:
- Maak generieke tekst concreet met domein, sector, senioriteit, klanttype of inhoudelijke context.
- Eisen moeten benoemen ervaring waarmee.
- Concurrenten moeten echte bedrijfsnamen zijn.
- No-go sourcing blijft exact de bedrijven uit feitenextractie.
- Eén bullet = één onderwerp. Lijsten voor taken/eisen/voorkeuren/USP/pullfactoren/arbeidsvoorwaarden bevatten precies 3 items.
- Eisen en voorkeuren zijn compact: maximaal 10 woorden per bullet, zonder opvulwoorden zoals 'relevante ervaring'.

FEITENEXTRACTIE:
{json.dumps(facts, ensure_ascii=False, indent=2)}

RESEARCH:
{json.dumps(research, ensure_ascii=False, indent=2)}

HUIDIGE JSON:
{json.dumps(data, ensure_ascii=False, indent=2)}
""".strip()




def build_presentation_prompt(data: Dict[str, Any], facts: Dict[str, Any], research: Dict[str, Any]) -> str:
    return f"""
Je bent presentation editor voor een Cooble startdocument. Verbeter alleen de presentatiekwaliteit.
Behoud de JSON-structuur exact en voeg geen feiten toe die niet in feitenextractie of research staan.

Harde regels:
- intake_samenvatting wordt één lopende tekst van 80-110 woorden. Geen bullets, geen kopjes, geen losse opsomming.
- taken_verantwoordelijkheden bevat exact 3 bullets.
- eisen bevat exact 3 bullets.
- voorkeuren bevat exact 3 bullets.
- usp_functie bevat exact 3 bullets.
- pullfactoren bevat exact 3 bullets.
- belangrijkste_arbeidsvoorwaarden bevat exact 3 bullets.
- Eén bullet = één onderwerp.
- Iedere bullet is concreet voor deze vacature of doelgroep.
- Geen generieke termen zoals relevante ervaring, passende kandidaat, dynamische omgeving.

FEITENEXTRACTIE:
{json.dumps(facts, ensure_ascii=False, indent=2)}

RESEARCH:
{json.dumps(research, ensure_ascii=False, indent=2)}

HUIDIGE STARTDOCUMENT JSON:
{json.dumps(data, ensure_ascii=False, indent=2)}
""".strip()

def generate_with_openai_pipeline(vacature: str, intake: str, linkedin_size: str, extra: str, status=None) -> Dict[str, Any]:
    """v0.7: meerstaps AI-pipeline met aparte presentation layer."""
    if status:
        status.write("Stap 1/5: feiten uit vacature en intake halen")
    facts = call_openai_json(build_fact_extraction_prompt(vacature, intake, extra), use_web=False)

    # Deterministische aanvulling voorkomt lege klant/functie/salaris als extractor iets mist.
    fallback = extract_basis_fallback(vacature, intake)
    for key_map in [("klantnaam", "klantnaam"), ("vacaturenaam", "vacaturenaam"), ("salaris", "salaris")]:
        fk, sk = key_map
        if is_empty_or_placeholder(facts.get(fk, "")) and fallback.get(sk):
            facts[fk] = fallback[sk]
    extracted_no_go = extract_no_go_companies_from_intake(intake + "\n" + extra)
    if extracted_no_go:
        merged = []
        for item in clean_list(facts.get("no_go_bedrijven", [])) + extracted_no_go:
            c = clean_company_name(item)
            if c and c not in merged:
                merged.append(c)
        facts["no_go_bedrijven"] = merged

    if status:
        status.write("Stap 2/5: arbeidsmarkt, doelgroep en concurrenten onderzoeken")
    research = call_openai_json(build_research_prompt(facts, linkedin_size, vacature, intake), use_web=True)

    if status:
        status.write("Stap 3/5: startdocument-content schrijven")
    data = call_openai_json(build_writer_prompt(facts, research, vacature, intake, linkedin_size, extra), use_web=False)

    if status:
        status.write("Stap 4/5: presentatiekwaliteit aanscherpen")
    try:
        data = call_openai_json(build_presentation_prompt(data, facts, research), use_web=False)
    except Exception as presentation_error:
        data.setdefault("kwaliteitscontrole", {}).setdefault("waarschuwingen", []).append(
            f"Presentation layer kon niet automatisch herschrijven: {presentation_error}"
        )

    if status:
        status.write("Stap 5/5: controleren op generieke tekst")
    data = apply_business_rules(data, intake + "\n" + extra, linkedin_size, vacature, extra)
    issues = collect_generic_issues(data)
    if issues:
        try:
            refined = call_openai_json(build_refine_prompt(data, facts, research, issues), use_web=False)
            data = apply_business_rules(refined, intake + "\n" + extra, linkedin_size, vacature, extra)
            data.setdefault("kwaliteitscontrole", {}).setdefault("waarschuwingen", []).append(
                "Automatische generieke-tekstcontrole uitgevoerd."
            )
        except Exception as refine_error:
            data.setdefault("kwaliteitscontrole", {}).setdefault("waarschuwingen", []).append(
                f"Generieke-tekstcontrole kon niet automatisch herschrijven: {refine_error}"
            )
    data.setdefault("kwaliteitscontrole", {})["pipeline"] = "v0.7: facts -> research -> writer -> presentation -> quality"
    return data


def generate_with_openai(prompt: str) -> Dict[str, Any]:
    """Compatibiliteitsfunctie voor oudere codepaden."""
    return call_openai_json(prompt, use_web=True)


def demo_data() -> Dict[str, Any]:
    return {
        "basisgegevens": {
            "klantnaam": "Voorbeeldklant",
            "vacaturenaam": "Consultant Waterkwaliteit",
            "datum": date.today().strftime("%d-%m-%Y"),
            "locatie": "Nederland",
            "uren": "32-40 uur",
            "salaris": "Schaal 8/9/10",
        },
        "intake_samenvatting": "Voor deze opdracht zoeken we een consultant die klanten kan adviseren over waterkwaliteit, vergunningen en industriële waterstromen. De nadruk ligt op iemand die technische of milieukundige kennis kan vertalen naar haalbare oplossingen voor bedrijven, waterschappen en drinkwaterorganisaties. Tijdens de intake is vooral benoemd dat de kandidaat breed naar water moet kunnen kijken, maar niet uit de hoek van hydrologie, ecologie of waterkwantiteit hoeft te komen. Een achtergrond in procestechnologie, werktuigbouwkunde of milieukunde kan juist interessant zijn wanneer de kandidaat adviesvaardig is en graag samenwerkt in multidisciplinaire projecten.",
        "functieprofiel": {
            "taken_verantwoordelijkheden": ["Adviseren over industriële waterkwaliteit", "Meedenken over vergunningen en compliance", "Ondersteunen van multidisciplinaire waterprojecten"],
            "usp_functie": ["Bouwen aan een groeiend waterteam", "Impact op duurzaam industrieel watergebruik", "Veel ruimte voor inhoudelijke ontwikkeling"],
        },
        "kandidaatprofiel": {
            "eisen": ["HBO werk- en denkniveau", "Ervaring met waterkwaliteit of afvalwater", "Adviesvaardigheid richting klanten"],
            "voorkeuren": ["Ervaring met vergunningstrajecten", "Achtergrond in procestechnologie", "Ervaring binnen industriële projecten"],
            "no_go_sourcing": [],
        },
        "voorwaarden": {"belangrijkste_arbeidsvoorwaarden": ["Hybride werken", "Ontwikkelmogelijkheden", "Goede pensioenregeling"]},
        "doelgroepanalyse": {
            "doelgroep_titel": "Voorbeeldfunctie",
            "verwachte_doelgroepgrootte": "± 500",
            "regio": "Nederland",
            "pullfactoren": ["Hybride werken", "Inhoudelijke complexiteit", "Autonomie"],
            "geslacht": {"man": "60%", "vrouw": "40%"},
            "leeftijdsverdeling": ["25-34: 30%", "35-44: 40%", "45-54: 20%", "55+: 10%"],
        },
        "sourcingplan": {
            "doelgroep": "Waterkwaliteitsadviseurs, milieukundig consultants en procestechnologen met ervaring in industriële waterstromen, vergunningen of afvalwaterprojecten.",
            "strategie": "Doelgroepgedreven sourcing met focus op vergelijkbare functies en organisaties.",
            "belangrijkste_functietitels": ["Consultant", "Adviseur", "Specialist"],
            "zoekrichting": ["LinkedIn sourcing", "Concurrenten op bedrijfsniveau", "Brede functietitelvarianten"],
            "toelichting": "Start breed en verfijn op inhoudelijke expertise en adviesvaardigheden.",
        },
        "concurrentenanalyse": {"relevant": True, "bedrijven": ["Witteveen+Bos", "Royal HaskoningDHV", "Sweco", "Antea Group", "Arcadis"], "toelichting": "Ingenieurs- en adviesbureaus waar vergelijkbare water- en milieuadviseurs werken."},
        "afspraken": ["Kandidaten worden voorgesteld na telefonische kennismaking.", "Feedback wordt zo snel mogelijk gedeeld.", "Bij profielwijzigingen wordt direct geschakeld."],
        "kwaliteitscontrole": {"ontbrekende_informatie": [], "aannames": [], "waarschuwingen": []},
    }




def validate_startdocument(data: Dict[str, Any]) -> Dict[str, List[str]]:
    """Geeft duidelijke kwaliteitsmeldingen terug voor de preview."""
    data = ensure_core_keys(data)
    errors: List[str] = []
    warnings: List[str] = []

    if is_empty_or_placeholder(get_nested(data, "basisgegevens.klantnaam", "")):
        errors.append("Klantnaam ontbreekt.")
    if is_empty_or_placeholder(get_nested(data, "basisgegevens.vacaturenaam", "")):
        errors.append("Vacaturenaam ontbreekt.")
    if len(str(data.get("intake_samenvatting", "")).split()) < 60:
        warnings.append("Intake-samenvatting is mogelijk te kort om de nuance goed over te brengen.")

    capped_lists = {
        "Taken & verantwoordelijkheden": get_nested(data, "functieprofiel.taken_verantwoordelijkheden", []),
        "Eisen": get_nested(data, "kandidaatprofiel.eisen", []),
        "Voorkeuren": get_nested(data, "kandidaatprofiel.voorkeuren", []),
        "USP's": get_nested(data, "functieprofiel.usp_functie", []),
        "Pullfactoren": get_nested(data, "doelgroepanalyse.pullfactoren", []),
        "Arbeidsvoorwaarden": get_nested(data, "voorwaarden.belangrijkste_arbeidsvoorwaarden", []),
    }
    for label, values in capped_lists.items():
        clean = clean_list(values)
        if len(clean) > 3:
            errors.append(f"{label} bevat meer dan 3 bullets.")
        if len(clean) < 3:
            warnings.append(f"{label} bevat minder dan 3 bullets.")

    generic_issues = collect_generic_issues(data)
    warnings.extend(generic_issues)

    competitors = clean_list(get_nested(data, "concurrentenanalyse.bedrijven", []))
    if not competitors:
        warnings.append("Concurrentenanalyse bevat nog geen bedrijven.")
    for company in competitors:
        if is_placeholder_company(company):
            errors.append("Concurrentenanalyse bevat een placeholder-bedrijf.")

    return {"errors": errors, "warnings": warnings}


def render_quality_check(data: Dict[str, Any]) -> None:
    result = validate_startdocument(data)
    errors = result["errors"]
    warnings = result["warnings"]
    if not errors and not warnings:
        st.success("Kwaliteitscheck geslaagd")
        return
    if errors:
        st.error("Kwaliteitscheck: actie nodig")
        for err in errors:
            st.write(f"- {err}")
    if warnings:
        st.warning("Aandachtspunten")
        for warn in warnings:
            st.write(f"- {warn}")


def editable_list(label: str, values: List[str], key: str, max_items: int = 6, *, hard_max: bool = False) -> List[str]:
    st.markdown(f"**{label}**")
    result = []
    values = values or []
    rows = max_items if hard_max else max(max_items, len(values))
    if hard_max and len(clean_list(values)) > max_items:
        st.caption(f"Let op: dit onderdeel is automatisch teruggebracht naar maximaal {max_items} bullets.")
    for i in range(rows):
        default = values[i] if i < len(values) else ""
        if key.startswith("afspraken"):
            val = st.text_area(f"{label} {i+1}", value=default, key=f"{key}_{i}", height=110, label_visibility="collapsed")
        else:
            val = st.text_input(f"{label} {i+1}", value=default, key=f"{key}_{i}", label_visibility="collapsed")
        if val.strip():
            result.append(val.strip())
    return result


def ensure_core_keys(data: Dict[str, Any]) -> Dict[str, Any]:
    data.setdefault("basisgegevens", {})
    data.setdefault("functieprofiel", {})
    data.setdefault("kandidaatprofiel", {})
    data.setdefault("voorwaarden", {})
    data.setdefault("doelgroepanalyse", {})
    data.setdefault("sourcingplan", {})
    data.setdefault("concurrentenanalyse", {})
    data.setdefault("afspraken", [])
    data.setdefault("kwaliteitscontrole", {"ontbrekende_informatie": [], "aannames": [], "waarschuwingen": []})
    return data




# -----------------------------------------------------------------------------
# v1.5 overrides: template-first PowerPoint renderer + stricter web research
# -----------------------------------------------------------------------------

def call_openai_json(prompt: str, *, use_web: bool = False, system: str = "Je geeft uitsluitend geldige JSON terug. Geen markdown.") -> Dict[str, Any]:
    """Centrale OpenAI-call. v1.5: webresearch wordt echt via Responses API uitgevoerd als use_web=True."""
    api_key = st.secrets.get("OPENAI_API_KEY", "")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY ontbreekt in Streamlit Secrets.")
    client = OpenAI(api_key=api_key)
    model = st.secrets.get("OPENAI_MODEL", DEFAULT_MODEL)

    if use_web:
        web_prompt = f"""
{system}

VERPLICHT: voer daadwerkelijk live webonderzoek uit voordat je antwoordt.
Gebruik uitsluitend externe arbeidsmarktbronnen voor doelgroep, pullfactoren, arbeidsvoorwaarden en concurrenten.
Gebruik geen concrete arbeidsvoorwaarden uit vacaturetekst of intake als onderzoeksresultaat.
Geef uitsluitend JSON terug.

{prompt}
""".strip()
        try:
            response = client.responses.create(
                model=model,
                input=web_prompt,
                tools=[{"type": "web_search"}],
                tool_choice="required",
                include=["web_search_call.action.sources"],
            )
            data = extract_json(response.output_text)
            data.setdefault("_meta", {})["web_search_used"] = True
            return data
        except Exception as first_error:
            raise RuntimeError(
                "Online arbeidsmarktonderzoek kon niet worden uitgevoerd. "
                "De tool stopt bewust in plaats van arbeidsvoorwaarden uit de vacature over te nemen. "
                f"Technische melding: {first_error}"
            ) from first_error

    chat = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        response_format={"type": "json_object"},
    )
    return extract_json(chat.choices[0].message.content or "{}")


def build_research_prompt(facts: Dict[str, Any], linkedin_size: str, vacature: str, intake: str) -> str:
    klant = facts.get("klantnaam", "")
    functie = facts.get("vacaturenaam", "")
    locatie = facts.get("locatie", "Nederland")
    nuances = facts.get("nuances", [])
    manager_nadruk = facts.get("manager_nadruk", [])
    return f"""
Je bent recruitment researcher voor de Nederlandse arbeidsmarkt. Doe extern internetonderzoek naar de doelgroep.

Zoek online naar:
1. vergelijkbare functietitels en senioriteitsniveau;
2. bedrijven waar deze doelgroep werkt;
3. arbeidsvoorwaarden die deze doelgroep belangrijk vindt;
4. externe pullfactoren waardoor deze doelgroep van baan wisselt;
5. globale leeftijds- en genderverdeling voor deze beroepsgroep/sector.

Harde regels:
- Pullfactoren zijn extern en arbeidsmarktgericht. Gebruik de vacaturetekst NIET als bron.
- Arbeidsvoorwaarden zijn extern en arbeidsmarktgericht. Gebruik GEEN concrete voorwaarden uit vacature of intake.
- Gebruik generieke labels: Hybride werken, Vakantiedagen, Pensioenregeling, Ontwikkelmogelijkheden, Mobiliteit.
- Concurrentenanalyse is altijd relevant en op bedrijfsniveau.
- Geef echte bedrijfsnamen. Nooit Bedrijf A/B/C, Concurrent 1 of Organisatie X.
- Doelgroepomschrijving moet specifiek zijn voor functie, domein, senioriteit en sector.
- Eén bullet = één onderwerp. Maximaal 3 pullfactoren en 3 arbeidsvoorwaarden.
- Als doelgroepgrootte uit LinkedIn is ingevuld, neem die letterlijk over.
- Voeg research_bronnen toe met korte bronlabels of domeinen die je hebt gebruikt.

Context om de juiste doelgroep te bepalen:
Klant: {klant}
Functie: {functie}
Locatie/regio: {locatie}
Doelgroepgrootte LinkedIn: {linkedin_size}
Nuances uit intake: {json.dumps(nuances, ensure_ascii=False)}
Manager nadruk: {json.dumps(manager_nadruk, ensure_ascii=False)}
No-go/check-eerst bedrijven: {json.dumps(facts.get('no_go_bedrijven', []), ensure_ascii=False)}

Geef uitsluitend JSON terug:
{{
  "doelgroep_titel": "",
  "doelgroep_omschrijving": "",
  "verwachte_doelgroepgrootte": "",
  "belangrijkste_functietitels": [],
  "pullfactoren": ["", "", ""],
  "belangrijkste_arbeidsvoorwaarden": ["", "", ""],
  "concurrenten_bedrijven": [],
  "zoekrichting": [],
  "geslacht": {{"man": "", "vrouw": ""}},
  "leeftijdsverdeling": ["25-34: %", "35-44: %", "45-54: %", "55+: %"],
  "research_bronnen": [],
  "research_toelichting": ""
}}
""".strip()


def build_employment_conditions_research_prompt(facts: Dict[str, Any], strict_retry: bool = False) -> str:
    functie = str(facts.get("vacaturenaam", "")).strip()
    retry = "" if not strict_retry else """
DIT IS EEN HERHALING OMDAT DE EERSTE UITKOMST GEEN DRIE GELDIGE ARBEIDSVOORWAARDEN OPLEVERDE.
Kies uitsluitend uit de toegestane lijst hieronder en controleer expliciet dat elk item een primaire of secundaire arbeidsvoorwaarde is.
"""
    allowed = json.dumps(ALLOWED_EMPLOYMENT_CONDITIONS, ensure_ascii=False)
    return f"""
Je bent arbeidsmarktonderzoeker. Doe ACTUEEL INTERNETONDERZOEK naar werknemersvoorkeuren binnen de Nederlandse arbeidsmarkt voor de beroepsgroep/functiefamilie:
{functie}
Land: Nederland

Onderzoek uitsluitend deze vraag:
WELKE DRIE PRIMAIRE OF SECUNDAIRE ARBEIDSVOORWAARDEN VINDT DEZE DOELGROEP HET BELANGRIJKST WANNEER ZIJ IN DIENST ZIJN?

{retry}
HARD REGELS:
- Gebruik verplicht live web_search en externe arbeidsmarktbronnen, werknemersenquêtes, brancheonderzoeken en doelgroepstudies.
- Gebruik de vacaturetekst, intake, werkgever, bedrijfsnaam en vacaturepagina NIET als bron.
- Een arbeidsvoorwaarde is een concrete contractuele/financiële/verlof-/mobiliteits-/werktijd-/opleidingsregeling.
- Werksfeer, cultuur, autonomie, impact, uitdaging, zingeving, doorgroei, carrièrekansen, ontwikkelmogelijkheden, werkzekerheid en inhoud van het werk zijn GEEN arbeidsvoorwaarden voor deze slide.
- Geef geen bedragen, percentages, aantallen dagen of werkgeversspecifieke invulling.
- Kies exact 3 verschillende categorieën uit deze toegestane lijst:
{allowed}
- Eén categorie per item, zonder toelichting in de array.
- Rangschik op relevantie voor deze beroepsgroep.
- Voeg de gebruikte bronnen/domeinen apart toe.

Geef uitsluitend JSON:
{{
  "belangrijkste_arbeidsvoorwaarden": ["", "", ""],
  "bronnen": [],
  "toelichting": ""
}}
""".strip()

def build_pullfactors_research_prompt(facts: Dict[str, Any], strict_retry: bool = False) -> str:
    functie = str(facts.get("vacaturenaam", "")).strip()
    retry_text = "" if not strict_retry else """
EXTRA CONTROLE BIJ DEZE HERHALING:
- Een eerdere uitkomst bevatte een ongeldige pullfactor.
- Verwijder arbeidsmarktkrapte, baankansen, tekort/schaarste, vacature-specifieke inhoud, bedrijfsnamen en werkgeverseigenschappen.
- Kies uitsluitend echte overstapmotieven/vacature-attractoren die uit extern doelgroep- of kandidatenonderzoek blijken.
"""
    return f"""
Je bent arbeidsmarktonderzoeker. Doe ACTUEEL INTERNETONDERZOEK naar de Nederlandse BEROEPSDOELGROEP achter deze functietitel:
Functietitel/functiefamilie: {functie}
Land: Nederland

Onderzoek uitsluitend deze vraag:
WELKE 3 FACTOREN BRENGEN DEZE BEROEPSDOELGROEP DAADWERKELIJK IN BEWEGING OM EEN ANDERE BAAN TE OVERWEGEN, EN WELKE NIET-FINANCIËLE ELEMENTEN WILLEN ZIJ DAAROM GRAAG TERUGZIEN IN EEN VACATURE?

BRONSCHEIDING:
- Gebruik VERPLICHT web_search en uitsluitend externe arbeidsmarkt-, kandidaten-, werknemers- en brancheonderzoeken.
- Gebruik GEEN vacaturetekst, intake, klantnaam, werkgever, taken, projecten, wetgeving, cultuurclaims of USP's van één organisatie.
- De functietitel/functiefamilie dient alleen om de juiste beroepsgroep te vinden.

WAT IS WEL EEN PULLFACTOR:
- een overstapmotief of vacature-attractor, zoals meer autonomie, inhoudelijke uitdaging, professionele ontwikkeling, doorgroeimogelijkheden, erkenning, zichtbare impact, strategische invloed of betere werk-privébalans — alleen wanneer onderzoek dit voor de doelgroep ondersteunt.

WAT IS GEEN PULLFACTOR:
- arbeidsmarktkansen, veel vacatures, schaarste, personeelstekort, baanzekerheid door krapte of hoge vraag;
- salaris, vakantiedagen, pensioen, leaseauto of andere concrete arbeidsvoorwaarden;
- kenmerken van de huidige vacature of werkgever;
- specifieke werkzaamheden/wetgeving/projecten uit één vacature.

Outputregels:
- Precies 3 pullfactoren.
- Elk item is één helder onderwerp in 2-7 natuurlijke woorden; maximaal 9 woorden.
- Geen samengestelde bullets.
- Formuleer als aantrekkingsfactor, niet als arbeidsmarktconstatering.
- Noem nooit een bedrijf of klant.
- Voeg gebruikte bronnen/domeinen toe.
{retry_text}
Geef uitsluitend JSON:
{{
  "pullfactoren": ["", "", ""],
  "bronnen": [],
  "toelichting": ""
}}
""".strip()


def pullfactors_contain_company(items: List[str], company: str) -> bool:
    company = re.sub(r"\\s+", " ", str(company or "")).strip().lower()
    if not company:
        return False
    company_tokens = [t for t in re.findall(r"[a-z0-9]+", company) if len(t) >= 4]
    for item in clean_list(items):
        low = item.lower()
        if company in low:
            return True
        # Ook herkenning op kenmerkende klanttokens, zodat bv. "Cosun Beet Company" wordt afgevangen.
        if any(tok in low for tok in company_tokens):
            return True
    return False


def pullfactors_are_invalid(items: List[str], company: str = "") -> bool:
    """Blokkeer uitkomsten die geen echte pullfactor zijn of vacature/werkgever-specifiek ogen."""
    forbidden = [
        r"arbeidsmarkt", r"baankans", r"baankansen", r"tekort", r"schaarste", r"veel vacatures",
        r"hoge vraag", r"gewild", r"krapte", r"personeelstekort", r"baanzekerheid door",
        r"salaris", r"pensioen", r"vakantiedag", r"leaseauto", r"eindejaarsuitkering",
    ]
    if pullfactors_contain_company(items, company):
        return True
    for item in clean_list(items):
        low = item.lower()
        if any(re.search(p, low) for p in forbidden):
            return True
    return len(clean_list(items)) != 3



def infer_occupation_family(facts: Dict[str, Any]) -> str:
    text = " ".join(str(facts.get(k, "")) for k in ["vacaturenaam", "doelgroep_titel", "locatie"]).lower()
    if re.search(r"\b(hvk|veiligheid|hse|qhse|safety|seveso|atex|procesindustrie|proces)\b", text):
        return "technical_safety_process"
    if re.search(r"\b(engineer|monteur|techniek|installatie|maintenance|werktuigbouw|elektro|bouw|civiel)\b", text):
        return "technical_engineering"
    if re.search(r"\b(it|software|developer|data|cloud|security|iam|architect)\b", text):
        return "it_digital"
    if re.search(r"\b(recruiter|hr|talent|people|personeel)\b", text):
        return "hr_recruitment"
    if re.search(r"\b(finance|controller|accountant|dbc|zorgcontrol|krediet|administratie)\b", text):
        return "finance_admin"
    if re.search(r"\b(zorg|verpleeg|kliniek|medisch|zorg)\b", text):
        return "healthcare"
    return "generic_professional"


def deterministic_demographics(facts: Dict[str, Any], research: Dict[str, Any]) -> Dict[str, Any]:
    """v2.4.4: stabiele demografie per functiefamilie.

    De webresearch bepaalt de afbakening/context, maar de presentatie gebruikt één
    vaste benchmark per beroepsfamilie. Daardoor wisselen man/vrouw en leeftijd
    niet meer per run voor dezelfde doelgroep.
    """
    family = infer_occupation_family(facts)
    benchmarks = {
        "technical_safety_process": {
            "geslacht": {"man": "80%", "vrouw": "20%"},
            "leeftijdsverdeling": ["15-24: 5%", "25-34: 20%", "35-49: 45%", "50+: 30%"],
        },
        "technical_engineering": {
            "geslacht": {"man": "85%", "vrouw": "15%"},
            "leeftijdsverdeling": ["15-24: 5%", "25-34: 25%", "35-49: 45%", "50+: 25%"],
        },
        "it_digital": {
            "geslacht": {"man": "80%", "vrouw": "20%"},
            "leeftijdsverdeling": ["15-24: 5%", "25-34: 40%", "35-49: 40%", "50+: 15%"],
        },
        "hr_recruitment": {
            "geslacht": {"man": "30%", "vrouw": "70%"},
            "leeftijdsverdeling": ["15-24: 5%", "25-34: 35%", "35-49: 45%", "50+: 15%"],
        },
        "finance_admin": {
            "geslacht": {"man": "50%", "vrouw": "50%"},
            "leeftijdsverdeling": ["15-24: 5%", "25-34: 30%", "35-49: 45%", "50+: 20%"],
        },
        "healthcare": {
            "geslacht": {"man": "20%", "vrouw": "80%"},
            "leeftijdsverdeling": ["15-24: 10%", "25-34: 30%", "35-49: 40%", "50+: 20%"],
        },
        "generic_professional": {
            "geslacht": {"man": "55%", "vrouw": "45%"},
            "leeftijdsverdeling": ["15-24: 5%", "25-34: 30%", "35-49: 45%", "50+: 20%"],
        },
    }
    result = dict(benchmarks.get(family, benchmarks["generic_professional"]))
    result["afbakening"] = f"{family} — gestabiliseerde benchmark op basis van beroepsfamilie"
    result["webresearch_bronnen"] = clean_list((research or {}).get("research_bronnen", []))
    return result

def _strip_target_size(value: Any) -> str:
    text = str(value or "").strip()
    text = re.sub(r"^[±+\-\s]+", "", text)
    return text


def _list3(items: Any) -> List[str]:
    vals = clean_list(items if isinstance(items, list) else [])
    return (vals + ["", "", ""])[:3]


def _replace_tokens_in_text_frame(shape, replacements: Dict[str, str]) -> None:
    """Vervang tekst in bestaande runs. Geen clear(), geen nieuwe fonts, geen nieuwe paragrafen."""
    if not (hasattr(shape, "text_frame") and shape.has_text_frame):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            txt = run.text
            if not txt or "{{" not in txt:
                continue
            for token, value in replacements.items():
                if token in txt:
                    txt = txt.replace(token, str(value or ""))
            run.text = txt


def _render_template_shape_v23(slide, shape, data: Dict[str, Any], replacements: Dict[str, str]) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _render_template_shape_v23(slide, child, data, replacements)
        return
    if not (hasattr(shape, "text_frame") and shape.has_text_frame):
        return
    original = full_text(shape)
    if "{{leeftijdsverdeling}}" in original:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        # De placeholder is alleen een anker; hij wordt onzichtbaar gemaakt.
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.text = ""
        img = create_age_chart_image(get_nested_v12(data, "doelgroepanalyse.leeftijdsverdeling", []))
        slide.shapes.add_picture(img, left, top, width=width, height=height)
        return
    _replace_tokens_in_text_frame(shape, replacements)


def generate_pptx(data: Dict[str, Any]) -> bytes:
    """v2.3: vul uitsluitend placeholders in het goedgekeurde Cooble-template."""
    template_path = get_template_path()
    prs = Presentation(str(template_path))

    tasks = _list3(get_nested_v12(data, "functieprofiel.taken_verantwoordelijkheden", []))
    eisen = _list3(get_nested_v12(data, "kandidaatprofiel.eisen", []))
    voorkeuren = _list3(get_nested_v12(data, "kandidaatprofiel.voorkeuren", []))
    usps = _list3(get_nested_v12(data, "functieprofiel.usp_functie", []))
    pulls = _list3(normalize_pullfactors(get_nested_v12(data, "doelgroepanalyse.pullfactoren", [])))
    voorwaarden = _list3(normalize_conditions(get_nested_v12(data, "voorwaarden.belangrijkste_arbeidsvoorwaarden", [])))
    afspraken = (clean_list(data.get("afspraken", [])) + ["", "", ""])[:3]
    nogo = clean_list(get_nested_v12(data, "kandidaatprofiel.no_go_sourcing", []))

    replacements = {
        "{{vacaturenaam}}": get_nested_v12(data, "basisgegevens.vacaturenaam", ""),
        "{{datum}}": get_nested_v12(data, "basisgegevens.datum", "") or date.today().strftime("%d-%m-%Y"),
        "{{intake_samenvatting}}": presentation_summary(data.get("intake_samenvatting", "")),
        "{{taak_1}}": tasks[0], "{{taak_2}}": tasks[1], "{{taak_3}}": tasks[2],
        "{{eis_1}}": eisen[0], "{{eis_2}}": eisen[1], "{{eis_3}}": eisen[2],
        "{{voorkeur_1}}": voorkeuren[0], "{{voorkeur_2}}": voorkeuren[1], "{{voorkeur_3}}": voorkeuren[2],
        "{{usp_1}}": usps[0], "{{usp_2}}": usps[1], "{{usp_3}}": usps[2],
        "{{no_go_sourcing}}": ", ".join(nogo),
        "{{salaris}}": normalize_salary_display(get_nested_v12(data, "basisgegevens.salaris", "")),
        "{{locatie}}": get_nested_v12(data, "basisgegevens.locatie", ""),
        "{{uren}}": get_nested_v12(data, "basisgegevens.uren", ""),
        "{{doelgroepgrootte}}": _strip_target_size(get_nested_v12(data, "doelgroepanalyse.verwachte_doelgroepgrootte", "")),
        "{{pull_1}}": pulls[0], "{{pull_2}}": pulls[1], "{{pull_3}}": pulls[2],
        "{{voorwaarde_1}}": voorwaarden[0], "{{voorwaarde_2}}": voorwaarden[1], "{{voorwaarde_3}}": voorwaarden[2],
        "{{geslacht_man}}": get_nested_v12(data, "doelgroepanalyse.geslacht.man", ""),
        "{{geslacht_vrouw}}": get_nested_v12(data, "doelgroepanalyse.geslacht.vrouw", ""),
        "{{afspraak_1}}": afspraken[0], "{{afspraak_2}}": afspraken[1], "{{afspraak_3}}": afspraken[2],
    }

    for slide in prs.slides:
        for shape in list(slide.shapes):
            _render_template_shape_v23(slide, shape, data, replacements)

    # Safety check: geen zichtbare placeholdertokens in het eindbestand.
    leftovers = []
    def collect(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                collect(sh.shapes)
            elif hasattr(sh, "text") and "{{" in sh.text:
                leftovers.append(sh.text)
    for slide in prs.slides:
        collect(slide.shapes)
    if leftovers:
        raise RuntimeError("Niet alle PowerPoint-placeholders konden worden ingevuld: " + " | ".join(leftovers[:5]))

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        result = Path(tmp.name).read_bytes()
    return result


# -----------------------------------------------------------------------------
# v2.4.2 overrides: strikte externe pullfactoren + vaste bewerkbare afspraken
# -----------------------------------------------------------------------------

DEFAULT_AFSPRAKEN = [
    "XX belt de kandidaten en doet de eerste screening. Mochten de kandidaten passen, worden ze via XX voorgesteld.",
    "Bij het afwijzen van kandidaten voorziet XX hen van feedback die ze terugkrijgen van de afdeling. Mochten kandidaten uitgenodigd worden voor een gesprek, neemt XX ook contact met hen op om kenbaar te maken dat er een interview met hen in wordt gepland. XX neemt het contact vanaf dit moment over.",
    "Bij wijzigingen wordt er geschakeld tussen XX, XX en XX.",
]

# Dit is bewust een gesloten woordenlijst. Webresearch bepaalt de RANGORDE,
# maar de uiteindelijke labels blijven algemene overstapmotieven en kunnen
# daardoor nooit per ongeluk vacature- of klantcontext bevatten.
ALLOWED_PULLFACTORS_V242 = [
    "Professionele ontwikkeling",
    "Werk-privébalans",
    "Autonomie",
    "Inhoudelijke uitdaging",
    "Doorgroeimogelijkheden",
    "Erkenning van expertise",
    "Flexibiliteit",
    "Werkzekerheid",
    "Maatschappelijke relevantie",
    "Leiderschap en cultuur",
]


def ensure_core_keys(data: Dict[str, Any]) -> Dict[str, Any]:
    data.setdefault("basisgegevens", {})
    data.setdefault("functieprofiel", {})
    data.setdefault("kandidaatprofiel", {})
    data.setdefault("voorwaarden", {})
    data.setdefault("doelgroepanalyse", {})
    data.setdefault("sourcingplan", {})
    data.setdefault("concurrentenanalyse", {})
    if not clean_list(data.get("afspraken", [])):
        data["afspraken"] = DEFAULT_AFSPRAKEN.copy()
    data.setdefault("kwaliteitscontrole", {"ontbrekende_informatie": [], "aannames": [], "waarschuwingen": []})
    return data


def build_pullfactors_research_prompt(facts: Dict[str, Any], strict_retry: bool = False) -> str:
    doelgroep = public_occupation_query(facts)
    retry = "" if not strict_retry else """
HERHAALCONTROLE:
- De vorige uitkomst voldeed niet aan de gesloten woordenlijst.
- Selecteer opnieuw precies drie labels uit de toegestane lijst.
"""
    return f"""
Je bent onafhankelijk arbeidsmarktonderzoeker. Doe ACTUEEL INTERNETONDERZOEK naar uitsluitend deze beroepsdoelgroep:
Doelgroep: {doelgroep}
Land: Nederland

ONDERZOEKSVRAAG:
Welke algemene factoren brengen professionals in deze beroepsgroep aantoonbaar in beweging om een andere baan te overwegen? Onderzoek werknemer-/kandidaatonderzoeken, arbeidsmarktstudies en brancheonderzoek. De vacature, werkgever en intake bestaan voor deze opdracht niet.

HARD BRONSCHEIDING:
- Gebruik verplicht web_search.
- Gebruik uitsluitend externe arbeidsmarkt-, kandidaat-, werknemer- en brancheonderzoeken.
- Gebruik GEEN vacaturetekst, vacaturepagina's, klantnaam, werkgever, intake, taken, projecten, locatie, wetgeving, bedrijfscultuur of USP's.
- Arbeidsmarktkrapte, baankansen, tekorten en schaarste zijn GEEN pullfactoren.
- Concrete arbeidsvoorwaarden (salaris, pensioen, vakantiedagen, bonus, leaseauto) horen in een andere module en zijn hier GEEN pullfactoren.

Selecteer op basis van het onderzoek precies 3 van deze algemene labels, in volgorde van relevantie voor de doelgroep:
{json.dumps(ALLOWED_PULLFACTORS_V242, ensure_ascii=False)}

BELANGRIJK:
- De drie waarden in pullfactoren moeten LETTERLIJK uit bovenstaande lijst komen.
- Geen eigen formuleringen en geen toelichtende zinnen in de pullfactoren-array.
- Geen sectorspecifieke/vacaturespecifieke formuleringen zoals 'complexe veiligheidsvraagstukken', 'impact op locatie', 'veiligheidscultuur' of 'organisatie in beweging'.
- Bronnen mogen wel apart worden opgenomen.
{retry}

Geef uitsluitend JSON:
{{
  "pullfactoren": ["", "", ""],
  "bronnen": [],
  "toelichting": ""
}}
""".strip()


def normalize_pullfactor_label(text: str) -> str:
    text = re.sub(r"\s+", " ", str(text or "").strip(" •-\n\t.,;:"))
    if not text:
        return ""
    low = text.lower()
    aliases = {
        "ontwikkeling": "Professionele ontwikkeling",
        "professionele groei": "Professionele ontwikkeling",
        "leren en ontwikkelen": "Professionele ontwikkeling",
        "work-life balance": "Werk-privébalans",
        "werk prive balans": "Werk-privébalans",
        "werk-privé balans": "Werk-privébalans",
        "vrijheid": "Autonomie",
        "eigenaarschap": "Autonomie",
        "uitdaging": "Inhoudelijke uitdaging",
        "doorgroei": "Doorgroeimogelijkheden",
        "erkenning": "Erkenning van expertise",
        "waardering": "Erkenning van expertise",
        "stabiliteit": "Werkzekerheid",
        "baanzekerheid": "Werkzekerheid",
        "sectorzekerheid": "Werkzekerheid",
        "zingeving": "Maatschappelijke relevantie",
        "cultuur": "Leiderschap en cultuur",
        "leiderschap": "Leiderschap en cultuur",
    }
    if text in ALLOWED_PULLFACTORS_V242:
        return text
    if low in aliases:
        return aliases[low]
    for key, value in aliases.items():
        if key in low:
            return value
    for allowed in ALLOWED_PULLFACTORS_V242:
        if allowed.lower() in low:
            return allowed
    # Geen vrije/contextuele formulering doorlaten.
    return ""


def normalize_pullfactors(items: List[str]) -> List[str]:
    out: List[str] = []
    for item in clean_list(items):
        label = normalize_pullfactor_label(item)
        if label and label not in out:
            out.append(label)
    # Alleen bij technisch onvolledige modeloutput aanvullen met neutrale algemene factoren.
    for fallback in ["Professionele ontwikkeling", "Werk-privébalans", "Autonomie"]:
        if len(out) >= 3:
            break
        if fallback not in out:
            out.append(fallback)
    return out[:3]


def pullfactors_are_invalid(items: List[str], company: str = "") -> bool:
    cleaned = clean_list(items)
    if len(cleaned) != 3:
        return True
    if pullfactors_contain_company(cleaned, company):
        return True
    # Model moet al letterlijk uit de gesloten lijst kiezen.
    if any(item not in ALLOWED_PULLFACTORS_V242 for item in cleaned):
        return True
    return len(set(cleaned)) != 3




def _occupation_research_context(facts: Dict[str, Any]) -> Dict[str, str]:
    """Maak een minimale, werkgever-neutrale context voor extern doelgroep onderzoek."""
    title = str(facts.get("vacaturenaam", "") or "").strip()
    # Klantnaam, vacature-USP's, interne nuances en concrete arbeidsvoorwaarden worden bewust niet doorgegeven.
    return {
        "functietitel": title,
        "locatie": str(facts.get("locatie", "Nederland") or "Nederland").strip(),
    }


def build_target_market_research_prompt(facts: Dict[str, Any], linkedin_size: str) -> str:
    ctx = _occupation_research_context(facts)
    return f"""
Je bent arbeidsmarktonderzoeker voor de Nederlandse recruitmentmarkt.
Voer VERPLICHT online onderzoek uit naar de BEROEPSGROEP rond onderstaande functietitel.

Doel: bepaal wie de relevante doelgroep is en bij welke werkgevers vergelijkbare professionals werken.

Strikte regels:
- Gebruik uitsluitend de functietitel/beroepsgroep en eventueel de geografische markt als onderzoekscontext.
- Gebruik GEEN werkgever, vacaturetekst, intake, vacature-USP's, interne organisatiecontext of arbeidsvoorwaarden als bron of context.
- Beschrijf de doelgroep beroepsmatig en specifiek: functiefamilie, specialisme, senioriteit, sectoren en gangbare alternatieve functietitels.
- Concurrenten zijn echte BEDRIJVEN/ORGANISATIES waar deze beroepsgroep aantoonbaar werkt; geef geen placeholders.
- Als LinkedIn-doelgroepgrootte is ingevuld, neem die exact over. Anders geef een conservatieve afgeronde indicatie op basis van online bronnen en benoem dat als schatting.
- Vermijd werkgever-specifieke termen in doelgroepomschrijving, functietitels en zoekrichting.

Onderzoekscontext:
Functietitel: {ctx['functietitel']}
Geografische markt: Nederland
Doelgroepgrootte gevonden op LinkedIn: {linkedin_size or 'niet ingevuld'}

Geef uitsluitend JSON:
{{
  "doelgroep_titel": "",
  "doelgroep_omschrijving": "",
  "verwachte_doelgroepgrootte": "",
  "belangrijkste_functietitels": [],
  "concurrenten_bedrijven": [],
  "zoekrichting": [],
  "research_toelichting": "",
  "bronnen": []
}}
""".strip()


def build_demographics_research_prompt(facts: Dict[str, Any]) -> str:
    ctx = _occupation_research_context(facts)
    return f"""
Je bent arbeidsmarktonderzoeker. Voer VERPLICHT online onderzoek uit naar de demografische samenstelling van de Nederlandse BEROEPSGROEP rond onderstaande functietitel.

Onderzoekscontext:
Functietitel: {ctx['functietitel']}
Land: Nederland

Strikte regels:
- Gebruik GEEN werkgever, vacaturetekst, intake, vacature-USP's of interne organisatiecontext.
- Zoek eerst naar zo direct mogelijke beroepsgroepdata. Geef voorkeur aan CBS, UWV, ROA, O*NET/ESCO-achtige beroepsindelingen met Nederlandse data, branche-/beroepsverenigingen en betrouwbare arbeidsmarktplatforms.
- Als er geen exacte functietiteldata bestaat, verbreed één stap naar de meest passende functiefamilie en leg dit kort uit in toelichting.
- Baseer man/vrouw en leeftijd op dezelfde of inhoudelijk vergelijkbare beroepsafbakening.
- Rond man/vrouw af op hele procenten die samen 100% vormen.
- Gebruik vaste leeftijdscategorieën: 15-24, 25-34, 35-49, 50+; percentages moeten samen circa 100% zijn.
- Geen vrije gok per run: gebruik gevonden broncijfers en een reproduceerbare afbakening.

Geef uitsluitend JSON:
{{
  "geslacht": {{"man": "", "vrouw": ""}},
  "leeftijdsverdeling": [
    "15-24: %",
    "25-34: %",
    "35-49: %",
    "50+: %"
  ],
  "demografie_toelichting": "",
  "bronnen": []
}}
""".strip()


def merge_research_parts(*parts: Dict[str, Any]) -> Dict[str, Any]:
    """Voeg onafhankelijke researchresultaten samen zonder eerdere waarden stil te overschrijven."""
    merged: Dict[str, Any] = {}
    all_sources: List[Any] = []
    for part in parts:
        if not isinstance(part, dict):
            continue
        for key, value in part.items():
            if key == "bronnen":
                if isinstance(value, list):
                    all_sources.extend(value)
                elif value:
                    all_sources.append(value)
                continue
            # Elke researchmodule beheert eigen velden; bij dubbeling wint de latere expliciete module.
            if value not in (None, "", [], {}):
                merged[key] = value
    if all_sources:
        seen = []
        for src in all_sources:
            if src not in seen:
                seen.append(src)
        merged["bronnen"] = seen
    return merged

def generate_with_openai_pipeline(vacature: str, intake: str, linkedin_size: str, extra: str, status=None) -> Dict[str, Any]:
    if status:
        status.write("Stap 1/8: feiten uit vacature en intake halen")
    facts = call_openai_json(build_fact_extraction_prompt(vacature, intake, extra), use_web=False)

    fallback = extract_basis_fallback(vacature, intake)
    for fk in ["klantnaam", "vacaturenaam", "salaris"]:
        if is_empty_or_placeholder(facts.get(fk, "")) and fallback.get(fk):
            facts[fk] = fallback[fk]
    facts["salaris"] = normalize_salary_display(facts.get("salaris", ""))

    extracted_no_go = extract_no_go_companies_from_intake(intake + "\n" + extra)
    if extracted_no_go:
        merged = []
        for item in clean_list(facts.get("no_go_bedrijven", [])) + extracted_no_go:
            c = clean_company_name(item)
            if c and c not in merged:
                merged.append(c)
        facts["no_go_bedrijven"] = merged

    if status:
        status.write("Stap 2/8: doelgroep en concurrenten online onderzoeken")
    market = call_openai_json(build_target_market_research_prompt(facts, linkedin_size), use_web=True)

    if status:
        status.write("Stap 3/8: belangrijkste arbeidsvoorwaarden online onderzoeken")
    conditions = call_openai_json(build_employment_conditions_research_prompt(facts), use_web=True)
    normalized_conditions = normalize_conditions(conditions.get("belangrijkste_arbeidsvoorwaarden", []))
    if employment_conditions_are_invalid(normalized_conditions):
        conditions = call_openai_json(build_employment_conditions_research_prompt(facts, strict_retry=True), use_web=True)
        normalized_conditions = normalize_conditions(conditions.get("belangrijkste_arbeidsvoorwaarden", []))
    if employment_conditions_are_invalid(normalized_conditions):
        raise RuntimeError(
            "Online arbeidsvoorwaardenonderzoek leverde geen drie geldige primaire/secundaire arbeidsvoorwaarden op. "
            "De tool stopt bewust in plaats van cultuur-, inhouds- of vacaturekenmerken te tonen."
        )
    conditions["belangrijkste_arbeidsvoorwaarden"] = normalized_conditions

    if status:
        status.write("Stap 4/8: pullfactoren onafhankelijk online onderzoeken")
    pull = call_openai_json(build_pullfactors_research_prompt(facts), use_web=True)
    normalized_pull = normalize_pullfactors(pull.get("pullfactoren", []))
    if pullfactors_are_invalid(normalized_pull, facts.get("klantnaam", "")):
        pull = call_openai_json(build_pullfactors_research_prompt(facts, strict_retry=True), use_web=True)
        normalized_pull = normalize_pullfactors(pull.get("pullfactoren", []))
    if pullfactors_are_invalid(normalized_pull, facts.get("klantnaam", "")):
        raise RuntimeError("Online pullfactoronderzoek leverde geen drie geldige algemene overstapmotieven op.")
    pull["pullfactoren"] = normalized_pull

    if status:
        status.write("Stap 5/8: leeftijd en man-vrouwverhouding online onderzoeken")
    demographics = call_openai_json(build_demographics_research_prompt(facts), use_web=True)
    research = merge_research_parts(market, conditions, pull, demographics)

    if status:
        status.write("Stap 6/8: startdocument-content schrijven")
    data = call_openai_json(build_writer_prompt(facts, research, vacature, intake, linkedin_size, extra), use_web=False)

    if status:
        status.write("Stap 7/8: presentatiekwaliteit aanscherpen")
    try:
        data = call_openai_json(build_presentation_prompt(data, facts, research), use_web=False)
    except Exception as presentation_error:
        data.setdefault("kwaliteitscontrole", {}).setdefault("waarschuwingen", []).append(str(presentation_error))

    if status:
        status.write("Stap 8/8: business rules toepassen")
    data = apply_business_rules(data, intake + "\n" + extra, linkedin_size, vacature, extra)

    # Externe research is altijd leidend voor deze velden.
    data.setdefault("doelgroepanalyse", {})["pullfactoren"] = normalize_pullfactors(research.get("pullfactoren", []))
    data.setdefault("voorwaarden", {})["belangrijkste_arbeidsvoorwaarden"] = presentation_bullets(normalize_conditions(research.get("belangrijkste_arbeidsvoorwaarden", [])), 3)
    stable_demo = deterministic_demographics(facts, research)
    data.setdefault("doelgroepanalyse", {})["geslacht"] = stable_demo.get("geslacht", {"man": "", "vrouw": ""})
    data.setdefault("doelgroepanalyse", {})["leeftijdsverdeling"] = stable_demo.get("leeftijdsverdeling", normalize_age_distribution(research.get("leeftijdsverdeling", [])))
    data.setdefault("basisgegevens", {})["salaris"] = normalize_salary_display(data.get("basisgegevens", {}).get("salaris", ""))

    # Afspraken zijn bewust een vast, bewerkbaar startpunt uit het voorbeeldtemplate.
    data["afspraken"] = DEFAULT_AFSPRAKEN.copy()
    data.setdefault("kwaliteitscontrole", {})["pipeline"] = "v2.4.4: facts -> external market -> external conditions -> closed-list external pull factors -> demographics -> writer -> compact candidate bullets + stable demographics + larger agreement editors"
    return ensure_core_keys(data)


st.title("📄 Startdocument Generator")
st.caption("Upload de vacature en intake. Controleer de preview en download daarna een nette PowerPoint in de vaste Cooble-template.")

with st.sidebar:
    st.header("Status")
    mode = st.radio("Modus", ["AI-generatie", "Testmodus zonder API-key"], index=0)
    st.caption("Testmodus gebruikt voorbeelddata en controleert of de app en PowerPoint-export werken.")
    if st.secrets.get("OPENAI_API_KEY", ""):
        st.success("API-key gevonden")
    else:
        st.warning("Geen API-key gevonden")

st.subheader("1. Input")
col1, col2 = st.columns(2)
with col1:
    st.markdown("**Vacature**")
    vacature_file = st.file_uploader("Vacature uploaden", type=["docx", "pdf", "txt"], key="vac_file")
    vacature_paste = st.text_area("Of plak vacaturetekst", height=240)
with col2:
    st.markdown("**Intake**")
    intake_file = st.file_uploader("Intake uploaden", type=["docx", "pdf", "txt"], key="intake_file")
    intake_paste = st.text_area("Of plak intake-notities", height=240)

c1, c2 = st.columns([1, 2])
with c1:
    linkedin_size = st.text_input("Doelgroepgrootte gevonden op LinkedIn", placeholder="bijv. ± 500")
with c2:
    extra_notes = st.text_area("Extra opmerkingen", placeholder="bijv. salaris niet benoemen / extra compact schrijven / regio belangrijk", height=80)

if "data" not in st.session_state:
    st.session_state.data = None

if st.button("Genereer analyse", type="primary"):
    try:
        vacature_text = (read_uploaded_file(vacature_file) + "\n" + vacature_paste).strip()
        intake_text = (read_uploaded_file(intake_file) + "\n" + intake_paste).strip()
        if mode == "AI-generatie" and not vacature_text:
            st.error("Voeg minimaal een vacaturetekst toe.")
        elif mode == "AI-generatie" and not intake_text:
            st.error("Voeg minimaal intake-informatie toe.")
        else:
            with st.status("Analyse wordt gemaakt...", expanded=True) as status:
                st.write("Vacature en intake uitlezen")
                if mode == "Testmodus zonder API-key":
                    st.write("Testdata laden")
                    data = demo_data()
                else:
                    st.write("AI-pipeline starten")
                    data = generate_with_openai_pipeline(vacature_text, intake_text, linkedin_size, extra_notes, status=status)
                data = apply_business_rules(data, intake_text + "\n" + extra_notes, linkedin_size, vacature_text, extra_notes)
                st.session_state.data = data
                status.update(label="Analyse klaar", state="complete", expanded=False)
            st.success("Analyse klaar. Controleer en pas eventueel aan.")
    except Exception as e:
        st.exception(e)

if st.session_state.data:
    data = ensure_core_keys(st.session_state.data)
    st.subheader("2. Preview & aanpassen")
    tabs = st.tabs(["Basis", "Functie", "Doelgroep", "Sourcing", "Afspraken"])

    with tabs[0]:
        b = data.setdefault("basisgegevens", {})
        c1, c2, c3 = st.columns(3)
        b["klantnaam"] = c1.text_input("Klantnaam", b.get("klantnaam", ""))
        b["vacaturenaam"] = c2.text_input("Vacaturenaam", b.get("vacaturenaam", ""))
        b["datum"] = date.today().strftime("%d-%m-%Y")
        c3.text_input("Datum", b["datum"], disabled=True, help="Altijd de generatiedatum van vandaag.")
        c4, c5, c6 = st.columns(3)
        b["locatie"] = c4.text_input("Locatie", b.get("locatie", ""))
        b["uren"] = c5.text_input("Uren", b.get("uren", ""))
        b["salaris"] = c6.text_input("Salaris", b.get("salaris", ""))
        data["intake_samenvatting"] = st.text_area("Intake / vacaturesamenvatting", data.get("intake_samenvatting", ""), height=170)

    with tabs[1]:
        f = data.setdefault("functieprofiel", {})
        k = data.setdefault("kandidaatprofiel", {})
        f["taken_verantwoordelijkheden"] = editable_list("Taken & verantwoordelijkheden", f.get("taken_verantwoordelijkheden", []), "taken", 3, hard_max=True)
        k["eisen"] = editable_list("Eisen", k.get("eisen", []), "eisen", 3, hard_max=True)
        k["voorkeuren"] = editable_list("Voorkeuren", k.get("voorkeuren", []), "voorkeuren", 3, hard_max=True)
        f["usp_functie"] = editable_list("USP's van de functie", f.get("usp_functie", []), "usp", 3, hard_max=True)
        k["no_go_sourcing"] = editable_list("No-go sourcing", k.get("no_go_sourcing", []), "nogo", 5)

    with tabs[2]:
        d = data.setdefault("doelgroepanalyse", {})
        v = data.setdefault("voorwaarden", {})
        c1, c2 = st.columns(2)
        d["verwachte_doelgroepgrootte"] = c1.text_input("Doelgroepgrootte", d.get("verwachte_doelgroepgrootte", ""))
        d["regio"] = c2.text_input("Regio", d.get("regio", "Nederland"))
        d["pullfactoren"] = editable_list("Pullfactoren", d.get("pullfactoren", []), "pull", 3, hard_max=True)
        v["belangrijkste_arbeidsvoorwaarden"] = editable_list("Belangrijkste arbeidsvoorwaarden", v.get("belangrijkste_arbeidsvoorwaarden", []), "av", 3, hard_max=True)
        g = d.setdefault("geslacht", {})
        c3, c4 = st.columns(2)
        g["man"] = c3.text_input("Geslacht man", g.get("man", ""))
        g["vrouw"] = c4.text_input("Geslacht vrouw", g.get("vrouw", ""))
        d["leeftijdsverdeling"] = editable_list("Leeftijdsverdeling", d.get("leeftijdsverdeling", []), "age", 4, hard_max=True)

    with tabs[3]:
        s = data.setdefault("sourcingplan", {})
        c = data.setdefault("concurrentenanalyse", {})
        s["doelgroep"] = st.text_area("Doelgroep", s.get("doelgroep", ""), height=100)
        s["strategie"] = st.text_area("Strategie", s.get("strategie", ""), height=120)
        s["zoekrichting"] = editable_list("Zoekrichting", s.get("zoekrichting", []), "zoek", 5)
        c["bedrijven"] = editable_list("Concurrenten / bedrijven", c.get("bedrijven", []), "conc", 8)
        c["toelichting"] = st.text_area("Toelichting concurrentenanalyse", c.get("toelichting", ""), height=90)

    with tabs[4]:
        data["afspraken"] = editable_list("Afspraken", data.get("afspraken", []), "afspraken", 5)

    st.session_state.data = data

    st.subheader("3. PowerPoint maken")
    try:
        data = apply_business_rules(data, "", linkedin_size)
        st.session_state.data = data
        pptx_bytes = generate_pptx(data)
        klant = get_nested(data, "basisgegevens.klantnaam", "klant") or "klant"
        vacature = get_nested(data, "basisgegevens.vacaturenaam", "vacature") or "vacature"
        filename = f"Startdocument_{klant}_{vacature}.pptx"
        filename = re.sub(r"[^A-Za-z0-9_\-\.]+", "_", filename)
        st.download_button(
            "Maak & download PowerPoint",
            pptx_bytes,
            filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
        )
    except Exception as e:
        st.exception(e)
